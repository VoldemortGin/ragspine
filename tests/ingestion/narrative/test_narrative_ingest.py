"""叙事批量入库编排测试（TDD 红色阶段）。

只验证外部行为：文件名 period 启发式（命中 / 不命中 / 歧义）、元数据两层来源
（显式 per-doc 映射优先、topic/entity 绝不猜测）、元数据透传到 chunk、
hash 比对幂等重入跳过、内容变化后版本递增、dry-run 零落库、no_text / failed
状态、文件夹扫描过滤、CLI（--db / --meta / --dry-run）。
fixture 在测试内现造（python-pptx / reportlab），sqlite 走 conftest 临时库。

红色预期：ragspine.ingestion.narrative.narrative_ingest / scripts.ingest_narrative 尚不存在，import 失败 = 红。
"""

import json
import os

import pytest
import rootutils

ROOT_DIR = rootutils.setup_root(os.getcwd(), indicator=".project-root", pythonpath=True)

from pptx import Presentation
from pptx.util import Inches
from reportlab.pdfgen.canvas import Canvas

from scripts.ingest_narrative import main as ingest_cli_main
from ragspine.retrieval.chunking.chunk_store import ChunkStore
from ragspine.ingestion.narrative.narrative_ingest import (
    STATUS_FAILED,
    STATUS_INGESTED,
    STATUS_NO_TEXT,
    STATUS_SKIPPED,
    ingest_narrative,
    period_from_filename,
)


# ---------------------------------------------------------------------------
# fixture 构造（测试内现造，确定性）
# ---------------------------------------------------------------------------

def _make_deck(path, paragraphs: list[str], notes: str | None = None) -> None:
    """单页 deck：每个段落一个文本框，可选演讲者备注。"""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    for i, text in enumerate(paragraphs):
        tb = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.5 + i * 0.9), Inches(8), Inches(0.8)
        )
        tb.text_frame.text = text
    if notes:
        slide.notes_slide.notes_text_frame.text = notes
    prs.save(str(path))


def _make_pdf(path, pages: list[str | None]) -> None:
    """按页造 PDF：字符串 = 文本页；None = 无文本层的扫描页形态。"""
    c = Canvas(str(path), pagesize=(595, 842))
    for content in pages:
        if content is None:
            c.setFillColorRGB(0.6, 0.6, 0.6)
            c.rect(40, 40, 500, 760, fill=1, stroke=0)
        else:
            c.setFont("Helvetica", 12)
            c.drawString(60, 700, content)
        c.showPage()
    c.save()


@pytest.fixture
def store(tmp_db_path):
    s = ChunkStore(tmp_db_path)
    s.init_schema()
    yield s
    s.close()


# ===========================================================================
# 文件名 period 启发式
# ===========================================================================

@pytest.mark.parametrize("name, expected", [
    ("ACME_HK_FY2024_review.pptx", "2024"),      # FY 模式 -> glossary 规范形
    ("group_results_2025H1.pdf", "2025H1"),     # 半年模式
    ("board_pack_2025Q1.pptx", "2025Q1"),       # 季度模式
    ("update fy2024 final.pptx", "2024"),       # 大小写不敏感
    ("FY2024H1_digest.pptx", "2024H1"),         # FY 前缀 + 半年后缀
])
def test_period_from_filename_hit(name, expected):
    assert period_from_filename(name) == expected


@pytest.mark.parametrize("name", [
    "meeting_notes.pptx",            # 无期间线索
    "2026-06-11 minutes.pdf",        # 裸年份 / 日期不算期间（绝不猜测）
    "FY2024_vs_FY2023_compare.pptx", # 多个不同期间 -> 歧义留空
])
def test_period_from_filename_miss(name):
    assert period_from_filename(name) == ""


# ===========================================================================
# 元数据策略 + 透传到 chunk
# ===========================================================================

def test_explicit_meta_passthrough_to_chunks(tmp_path, store):
    """显式 per-doc 元数据（含 valid_as_of）逐字段透传到落库的块。"""
    deck = tmp_path / "qbr_FY2024.pptx"
    _make_deck(deck, ["CPL loss attribution detail."], notes="Follow-up pending.")
    meta = {"qbr_FY2024.pptx": {
        "topic": "QBR",
        "entity": "ACME_CN",
        "geography": "CN",
        "language": "en",
        "sensitivity": "RESTRICTED",
        "valid_as_of": "2026-06-01",
    }}

    report = ingest_narrative([deck], store, meta_by_doc=meta)

    assert report.dry_run is False
    assert len(report.files) == 1
    f = report.files[0]
    assert f.status == STATUS_INGESTED
    assert f.doc_id == "qbr_FY2024.pptx"
    assert f.n_chunks == store.count() > 0

    rows = store.iter_chunks(doc_id="qbr_FY2024.pptx")
    assert all(r.topic == "QBR" for r in rows)
    assert all(r.entity == "ACME_CN" for r in rows)
    assert all(r.geography == "CN" for r in rows)
    assert all(r.language == "en" for r in rows)
    assert all(r.sensitivity == "RESTRICTED" for r in rows)
    assert all(r.valid_as_of == "2026-06-01" for r in rows)
    # meta 未给 period -> 文件名启发式补上
    assert all(r.period == "2024" for r in rows)


def test_no_meta_never_guesses_topic_entity(tmp_path, store):
    """缺省元数据：period 走文件名启发式，topic/entity 留空就是留空。"""
    deck = tmp_path / "townhall_2025H1.pptx"
    _make_deck(deck, ["Agency expansion remains on track."])

    ingest_narrative([deck], store)

    rows = store.iter_chunks(doc_id="townhall_2025H1.pptx")
    assert rows
    assert all(r.topic == "" for r in rows)
    assert all(r.entity == "" for r in rows)
    assert all(r.geography == "" for r in rows)
    assert all(r.period == "2025H1" for r in rows)


def test_explicit_period_overrides_filename(tmp_path, store):
    """显式 period 优先于文件名启发式。"""
    deck = tmp_path / "review_FY2024.pptx"
    _make_deck(deck, ["Half-year deep dive."])

    ingest_narrative([deck], store, meta_by_doc={"review_FY2024.pptx": {"period": "2025H1"}})

    rows = store.iter_chunks(doc_id="review_FY2024.pptx")
    assert rows and all(r.period == "2025H1" for r in rows)


def test_unknown_meta_key_rejected(tmp_path, store):
    """元数据映射里的未知字段（疑似笔误）直接 ValueError，不静默忽略。"""
    deck = tmp_path / "deck.pptx"
    _make_deck(deck, ["Some narrative."])
    with pytest.raises(ValueError):
        ingest_narrative([deck], store, meta_by_doc={"deck.pptx": {"topicx": "QBR"}})


# ===========================================================================
# 幂等重入 / 版本
# ===========================================================================

def test_reingest_unchanged_file_skipped(tmp_path, store):
    """同文件二次入库：hash 未变 -> skipped，库内不重复、版本不递增。"""
    deck = tmp_path / "deck.pptx"
    _make_deck(deck, ["Narrative paragraph one.", "Narrative paragraph two."])

    first = ingest_narrative([deck], store)
    assert first.files[0].status == STATUS_INGESTED
    n_after_first = store.count()

    second = ingest_narrative([deck], store)
    assert second.files[0].status == STATUS_SKIPPED
    assert store.count() == n_after_first
    assert {r.version for r in store.iter_chunks(doc_id="deck.pptx")} == {1}


def test_changed_file_reingested_with_new_version(tmp_path, store):
    """文件内容变化 -> 重新入库，走 chunk_store 版本化语义（活跃版本递增）。"""
    deck = tmp_path / "deck.pptx"
    _make_deck(deck, ["Original narrative."])
    ingest_narrative([deck], store)

    _make_deck(deck, ["Original narrative.", "Newly appended attribution."])
    report = ingest_narrative([deck], store)

    assert report.files[0].status == STATUS_INGESTED
    rows = store.iter_chunks(doc_id="deck.pptx")
    assert {r.version for r in rows} == {2}
    assert any("Newly appended attribution" in r.text for r in rows)


# ===========================================================================
# dry-run / no_text / failed / 文件夹扫描
# ===========================================================================

def test_dry_run_reports_but_writes_nothing(tmp_path, store):
    """dry-run：报告完整（将要入库的 chunk 数），但块库与登记台账零写入。"""
    deck = tmp_path / "deck.pptx"
    _make_deck(deck, ["Narrative for dry run."])

    report = ingest_narrative([deck], store, dry_run=True)
    assert report.dry_run is True
    assert report.files[0].status == STATUS_INGESTED
    assert report.files[0].n_chunks > 0
    assert store.count() == 0

    # 台账也未登记：随后真实入库不会被误判为 skipped。
    real = ingest_narrative([deck], store)
    assert real.files[0].status == STATUS_INGESTED
    assert store.count() > 0


def test_all_scanned_pdf_is_no_text(tmp_path, store):
    """全扫描 PDF：no_text 状态 + 跳过页计数，不落库。"""
    pdf = tmp_path / "scanned_minutes.pdf"
    _make_pdf(pdf, [None, None])

    report = ingest_narrative([pdf], store)
    f = report.files[0]
    assert f.status == STATUS_NO_TEXT
    assert f.n_skipped_pages == 2
    assert f.n_chunks == 0
    assert store.count() == 0


def test_failed_file_does_not_break_batch(tmp_path, store):
    """坏文件 failed 且记录原因，同批其它文件照常入库。"""
    folder = tmp_path / "docs"
    folder.mkdir()
    (folder / "broken.pptx").write_bytes(b"this is not a pptx")
    _make_deck(folder / "good.pptx", ["Healthy narrative."])

    report = ingest_narrative(folder, store)
    by_id = {f.doc_id: f for f in report.files}
    assert by_id["broken.pptx"].status == STATUS_FAILED
    assert by_id["broken.pptx"].error
    assert by_id["good.pptx"].status == STATUS_INGESTED
    assert store.count() > 0


def test_folder_scan_only_supported_types(tmp_path, store):
    """文件夹输入只收 pptx/pdf，忽略其它类型与 Office 临时文件。"""
    folder = tmp_path / "docs"
    folder.mkdir()
    _make_deck(folder / "deck.pptx", ["Narrative."])
    _make_pdf(folder / "report.pdf", ["Some report text."])
    (folder / "readme.txt").write_text("ignore me", encoding="utf-8")
    (folder / "~$deck.pptx").write_bytes(b"office lock file")

    report = ingest_narrative(folder, store)
    assert sorted(f.doc_id for f in report.files) == ["deck.pptx", "report.pdf"]
    assert all(f.status == STATUS_INGESTED for f in report.files)


# ===========================================================================
# CLI（scripts/ingest_narrative.py）
# ===========================================================================

def test_cli_dry_run_then_ingest(tmp_path):
    """CLI：--dry-run 零落库；正式跑落库且 --meta 元数据生效；坏文件时退出码 1。"""
    folder = tmp_path / "docs"
    folder.mkdir()
    _make_deck(folder / "qbr_FY2024.pptx", ["CPL attribution narrative."])
    db = tmp_path / "chunks.db"
    meta_file = tmp_path / "meta.json"
    meta_file.write_text(
        json.dumps({"qbr_FY2024.pptx": {"topic": "QBR", "entity": "ACME_CN"}}),
        encoding="utf-8",
    )

    rc = ingest_cli_main([str(folder), "--db", str(db), "--meta", str(meta_file), "--dry-run"])
    assert rc == 0
    s = ChunkStore(db)
    s.init_schema()
    assert s.count() == 0

    rc = ingest_cli_main([str(folder), "--db", str(db), "--meta", str(meta_file)])
    assert rc == 0
    rows = s.iter_chunks(doc_id="qbr_FY2024.pptx")
    assert rows and all(r.topic == "QBR" and r.entity == "ACME_CN" for r in rows)
    s.close()

    (folder / "broken.pptx").write_bytes(b"garbage")
    rc = ingest_cli_main([str(folder), "--db", str(db)])
    assert rc == 1
