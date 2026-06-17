"""叙事文本抽取测试（TDD 红色阶段）。

只验证外部行为：pptx 文本框（含多段落、空框）+ 演讲者备注的抽取与定位、
表格内容确实被跳过、数字型 PDF 逐页抽取与定位、无文本层页跳过计数告警、
to_text() 空行分段（喂 src/chunking.py 的契约）、后缀分发。
fixture 在测试内现造（python-pptx / reportlab），零静态文件、零网络。

红色预期：ragspine.ingestion.narrative.narrative_extract 尚不存在，import 失败 = 红。
"""

import os

import pytest
import rootutils

ROOT_DIR = rootutils.setup_root(os.getcwd(), indicator=".project-root", pythonpath=True)

from pptx import Presentation
from pptx.util import Inches
from reportlab.pdfgen.canvas import Canvas

from ragspine.ingestion.narrative.narrative_extract import (
    NarrativeDoc,
    NarrativeSegment,
    extract_narrative,
    extract_pdf_narrative,
    extract_pptx_narrative,
)

# 表格哨兵内容：绝不允许出现在叙事抽取结果里（表格数字归结构化通路）。
TABLE_MARKER = "TABLE_ONLY_MARKER_4321"
TABLE_NUMBER = "99999"


# ---------------------------------------------------------------------------
# fixture 构造（测试内现造，确定性）
# ---------------------------------------------------------------------------

def _make_deck(path) -> None:
    """两页 deck：文本框（含多段落 / 空框）+ 表格（哨兵）+ 演讲者备注。"""
    prs = Presentation()
    blank = prs.slide_layouts[6]

    s1 = prs.slides.add_slide(blank)
    tb1 = s1.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(8), Inches(0.8))
    tb1.text_frame.text = "FY2024 Hong Kong performance review"
    tb2 = s1.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(8), Inches(1.2))
    tb2.text_frame.text = "REVENUE grew strongly in 2024."
    p = tb2.text_frame.add_paragraph()
    p.text = "Growth was driven by   agency   expansion."  # 内部多空格，应折叠
    tbl = s1.shapes.add_table(2, 2, Inches(0.5), Inches(3.5), Inches(6), Inches(1)).table
    tbl.cell(0, 0).text = TABLE_MARKER
    tbl.cell(1, 1).text = TABLE_NUMBER

    s2 = prs.slides.add_slide(blank)
    # 空文本框放在最前：既不产段、也不占 frame 序号。
    s2.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(4), Inches(0.5))
    tb3 = s2.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(8), Inches(0.8))
    tb3.text_frame.text = "Regulatory outlook remains stable."
    s2.notes_slide.notes_text_frame.text = "Speaker note: CPL loss attribution pending."

    prs.save(str(path))


def _make_empty_deck(path) -> None:
    """一页空 deck：无任何文本框 / 表格 / 备注。"""
    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[6])
    prs.save(str(path))


def _make_pdf(path, pages: list[str | None]) -> None:
    """按页造 PDF：字符串 = 数字型文本页；None = 仅位图式矢量块、无文本层。"""
    c = Canvas(str(path), pagesize=(595, 842))
    for content in pages:
        if content is None:
            c.setFillColorRGB(0.6, 0.6, 0.6)
            c.rect(40, 40, 500, 760, fill=1, stroke=0)
        else:
            c.setFont("Helvetica", 12)
            c.drawString(60, 700, content)
        c.showPage()
    c.save()


@pytest.fixture
def deck_path(tmp_path):
    path = tmp_path / "deck_FY2024.pptx"
    _make_deck(path)
    return path


@pytest.fixture
def mixed_pdf_path(tmp_path):
    """3 页 PDF：第 1、3 页有文本层，第 2 页无文本层（扫描页形态）。"""
    path = tmp_path / "report.pdf"
    _make_pdf(path, [
        "Regulatory update for Hong Kong insurance market.",
        None,
        "The regulator issued new guidance on sales conduct.",
    ])
    return path


# ===========================================================================
# pptx：文本框 + 备注抽取与定位
# ===========================================================================

def test_pptx_segments_and_locators(deck_path):
    """文本框按 slide/frame 定位、备注 slide=N,notes；空框不产段不占序号。"""
    doc = extract_pptx_narrative(deck_path)
    assert isinstance(doc, NarrativeDoc)
    assert all(isinstance(s, NarrativeSegment) for s in doc.segments)
    assert [s.source_locator for s in doc.segments] == [
        "slide=1,frame=1",
        "slide=1,frame=2",
        "slide=2,frame=1",
        "slide=2,notes",
    ]
    assert doc.segments[0].text == "FY2024 Hong Kong performance review"
    # 多段落文本框：段落各占一行；内部连续空白折叠为单空格。
    assert doc.segments[1].text == (
        "REVENUE grew strongly in 2024.\nGrowth was driven by agency expansion."
    )
    assert doc.segments[2].text == "Regulatory outlook remains stable."
    assert doc.segments[3].text == "Speaker note: CPL loss attribution pending."


def test_pptx_table_content_skipped(deck_path):
    """表格内容（哨兵文本与数字）绝不进入叙事抽取结果。"""
    doc = extract_pptx_narrative(deck_path)
    full_text = doc.to_text()
    assert TABLE_MARKER not in full_text
    assert TABLE_NUMBER not in full_text


def test_pptx_doc_fields(deck_path):
    """doc_id=文件名、file_hash 为 sha256 十六进制串、pptx 无跳页。"""
    doc = extract_pptx_narrative(deck_path)
    assert doc.doc_id == "deck_FY2024.pptx"
    assert len(doc.file_hash) == 64 and all(c in "0123456789abcdef" for c in doc.file_hash)
    assert doc.skipped_pages == 0


def test_pptx_to_text_blank_line_separated(deck_path):
    """to_text() = 各段以空行连接（配合 chunking 的段落切分契约）。"""
    doc = extract_pptx_narrative(deck_path)
    assert doc.to_text() == "\n\n".join(s.text for s in doc.segments)
    assert "\n\n" in doc.to_text()


def test_pptx_empty_deck(tmp_path):
    """无文本内容的 deck -> 空段列表、空文本。"""
    path = tmp_path / "empty.pptx"
    _make_empty_deck(path)
    doc = extract_pptx_narrative(path)
    assert doc.segments == []
    assert doc.to_text() == ""


# ===========================================================================
# PDF：逐页抽取 + 无文本层页跳过计数
# ===========================================================================

def test_pdf_pages_extracted_with_locators(mixed_pdf_path):
    """有文本层的页逐页成段，locator 用真实页号（含被跳过页的占位）。"""
    doc = extract_pdf_narrative(mixed_pdf_path)
    assert [s.source_locator for s in doc.segments] == ["page=1", "page=3"]
    assert "Regulatory update" in doc.segments[0].text
    assert "guidance on sales conduct" in doc.segments[1].text


def test_pdf_scanned_page_skipped_and_counted(mixed_pdf_path):
    """无文本层的页跳过并计数，告警里点名页号。"""
    doc = extract_pdf_narrative(mixed_pdf_path)
    assert doc.skipped_pages == 1
    assert any("page=2" in w for w in doc.warnings)


def test_pdf_all_scanned(tmp_path):
    """全扫描 PDF：零段、跳页数=页数。"""
    path = tmp_path / "scanned.pdf"
    _make_pdf(path, [None, None])
    doc = extract_pdf_narrative(path)
    assert doc.segments == []
    assert doc.skipped_pages == 2


def test_pdf_doc_fields(mixed_pdf_path):
    doc = extract_pdf_narrative(mixed_pdf_path)
    assert doc.doc_id == "report.pdf"
    assert len(doc.file_hash) == 64


# ===========================================================================
# 后缀分发
# ===========================================================================

def test_dispatch_by_suffix(deck_path, mixed_pdf_path):
    """extract_narrative 按后缀分发，结果与专用入口一致。"""
    via_dispatch = extract_narrative(deck_path)
    direct = extract_pptx_narrative(deck_path)
    assert [s.source_locator for s in via_dispatch.segments] == [
        s.source_locator for s in direct.segments
    ]
    pdf_doc = extract_narrative(mixed_pdf_path)
    assert [s.source_locator for s in pdf_doc.segments] == ["page=1", "page=3"]


def test_dispatch_unsupported_suffix(tmp_path):
    path = tmp_path / "notes.txt"
    path.write_text("plain text", encoding="utf-8")
    with pytest.raises(ValueError):
        extract_narrative(path)
