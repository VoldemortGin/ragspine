"""统一多格式分发入口 ingest_file 的红色测试（TDD，3A 任务）。

接通入库主路径：数字通路目前生产只吃 Excel（ingestion.py 仅 import
xlsx_styled_extractor）；pptx / pdf 抽取器与 pdf_router 已实现但无生产调用方；
ingest_excel 收 queue 参数却从不 enqueue（n_enqueued_review 恒 0）。本批接通。

只验证外部行为：给定各格式 fixture（xlsx / pptx / 数字型 pdf / 扫描型 pdf）+
就地构造的刁钻输入（实体不可解析 / 有色无映射 / 不支持后缀），断言
ingest_file 的对外输出（IngestReport 内容 + fact_store / review_queue 状态）。

红色预期：所有用例因 `from ragspine.ingestion.structured.ingestion import ingest_file` ImportError
（该入口尚未实现）而 FAIL。import 放在每个测试体内首行，使其作为用例
FAILURE（而非 collection ERROR）暴露——沿用 tests/test_ingestion.py 的红色阶段约定。

每个用例 docstring 带中文 user story。
"""

import os

import pytest
import rootutils

ROOT_DIR = rootutils.setup_root(os.getcwd(), indicator=".project-root", pythonpath=True)

from openpyxl import Workbook
from openpyxl.styles import PatternFill
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches

from ragspine.extraction.color.color_semantics import ColorMapping, LegendEntry, MappingRegistry
from ragspine.storage.fact_store import FactStore, VISIBLE_REVIEW_STATUSES
from ragspine.ingestion.structured.ingestion import IngestReport, ingest_excel
from ragspine.ingestion.review.review_queue import ReviewQueue

# excel_styled_fixture.xlsx 的 HK_Performance sheet：A1='ACME Hong Kong'、
# 首行 FY2022..FY2024、首列 REVENUE/NEWSALES/PROFIT/ROE = 4 指标 × 3 期间 = 12 候选事实。
HK_SCOPE = "excel_styled_fixture.xlsx"
EXPECTED_HK_FACTS = 12

# 受控实体（home 公司 profile 可归一）：grid 的 A1 / sheet / 标题写它，实体可解析。
RESOLVABLE_ENTITY_TITLE = "ACME Hong Kong"   # -> ACME_HK
# 语义色（与 Excel 线同一套：黄色 = 新产品线）。
YELLOW = "FFFF00"


# --------------------------------------------------------------------------- #
# 三件套 fixture：只构造对象（开 sqlite 连接，不触发任何契约方法）。
# init_schema / 装配 active 映射等放进辅助函数，由测试体调用。
# --------------------------------------------------------------------------- #
@pytest.fixture
def store(tmp_sqlite_factory):
    fs = FactStore(tmp_sqlite_factory("facts"))
    yield fs
    fs.close()


@pytest.fixture
def queue(tmp_sqlite_factory):
    q = ReviewQueue(tmp_sqlite_factory("queue"))
    yield q
    q.close()


@pytest.fixture
def registry(tmp_sqlite_factory):
    reg = MappingRegistry(tmp_sqlite_factory("registry"))
    yield reg
    reg.close()


def _init_three(store, registry, queue) -> None:
    """建三库表（无颜色映射版，供不依赖颜色 tag 的用例用）。"""
    store.init_schema()
    queue.init_schema()
    registry.init_schema()


def _setup_hk_mapping(store, registry, queue, ground_truth) -> None:
    """建表并装配与 excel fixture HK_Performance 图例一致的 active 颜色映射。"""
    _init_three(store, registry, queue)
    legend = ground_truth["sheets"]["HK_Performance"]["legend_expect"]
    entries = [
        LegendEntry(
            rgb=e["rgb"],
            meaning=e["meaning"],
            tag_key=e["tag_key"],
            tag_value=e["tag_value"],
        )
        for e in legend
    ]
    version = registry.register_draft(ColorMapping(scope=HK_SCOPE, entries=entries))
    registry.confirm(HK_SCOPE, version, actor="sme_fin", note="fixture 图例")


# --------------------------------------------------------------------------- #
# 就地构造各格式 fixture（确定性、零外部敏感数据）。
# --------------------------------------------------------------------------- #
def _make_pptx_financial(path, *, entity_title: str, colored: bool = False) -> None:
    """造一张【实体可解析】的 pptx 财务表：R1C1=entity_title、首行期间、首列指标。

    colored=True 时给数值格上黄色填充（供「有色无映射」用例验证颜色 tag 需翻译）。
    """
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    # 表：2 行 × 2 列。R1C1=实体、R1C2=期间表头、R2C1=指标、R2C2=数值。
    tbl = slide.shapes.add_table(
        2, 2, Inches(0.5), Inches(1.0), Inches(8), Inches(2)
    ).table
    tbl.cell(0, 0).text = entity_title
    tbl.cell(0, 1).text = "FY2024"
    tbl.cell(1, 0).text = "REVENUE"
    tbl.cell(1, 1).text = "2680"
    if colored:
        tbl.cell(1, 1).fill.solid()
        tbl.cell(1, 1).fill.fore_color.rgb = RGBColor.from_string(YELLOW)
    prs.save(str(path))


def _make_xlsx_financial(
    path, *, sheet_title: str, a1: str, colored: bool = False
) -> None:
    """造一张 xlsx 财务表：A1=a1（实体源）、首行期间、首列指标、B2 数值。

    sheet_title / a1 决定实体能否归一；colored=True 时给 B2 上黄色填充。
    """
    wb = Workbook()
    ws = wb.active
    ws.title = sheet_title
    ws["A1"] = a1
    ws["B1"] = "FY2024"
    ws["A2"] = "REVENUE"
    ws["B2"] = 2680
    if colored:
        ws["B2"].fill = PatternFill(
            start_color="FFFFFF00", end_color="FFFFFF00", fill_type="solid"
        )
    wb.save(str(path))


def _fact_keys(rows) -> set[tuple]:
    """事实唯一键集合（与 fact_metric 唯一索引一致）。

    rows 为 store.execute_read 返回的 sqlite3.Row（按列名下标访问）。
    """
    return {
        (r["metric_code"], r["entity"], r["period_type"], r["period"], r["channel"])
        for r in rows
    }


# ===========================================================================
# D1 — xlsx 平价：ingest_file 与 ingest_excel 对同一 xlsx 入库结果一致
# ===========================================================================
def test_d1_xlsx_parity_fact_count(
    store, registry, queue, excel_fixture_path, ground_truth
):
    """user story：作为接手方，我要 ingest_file 走 xlsx 分支与既有 ingest_excel
    入库条数完全一致，确保统一入口不丢不重既有 Excel 能力。"""
    from ragspine.ingestion.structured.ingestion import ingest_file

    _setup_hk_mapping(store, registry, queue, ground_truth)
    report = ingest_file(excel_fixture_path, store, registry, queue)
    assert report.n_facts_ingested == EXPECTED_HK_FACTS
    assert store.count() == EXPECTED_HK_FACTS


def test_d1_xlsx_parity_same_keys(
    store, registry, queue, tmp_sqlite_factory, excel_fixture_path, ground_truth
):
    """user story：同一 xlsx 经 ingest_file 与 ingest_excel 入库的事实键集合一致
    （指标×实体×期间×渠道逐键相同），证明分发入口对 Excel 是行为等价封装。"""
    from ragspine.ingestion.structured.ingestion import ingest_file

    # 入口 A：ingest_file
    _setup_hk_mapping(store, registry, queue, ground_truth)
    ingest_file(excel_fixture_path, store, registry, queue)
    keys_file = _fact_keys(store.execute_read("SELECT * FROM fact_metric"))

    # 入口 B：ingest_excel（独立三库）
    store_b = FactStore(tmp_sqlite_factory("facts_b"))
    reg_b = MappingRegistry(tmp_sqlite_factory("registry_b"))
    q_b = ReviewQueue(tmp_sqlite_factory("queue_b"))
    try:
        _setup_hk_mapping(store_b, reg_b, q_b, ground_truth)
        ingest_excel(excel_fixture_path, store_b, reg_b, q_b)
        keys_excel = _fact_keys(store_b.execute_read("SELECT * FROM fact_metric"))
    finally:
        store_b.close()
        reg_b.close()
        q_b.close()

    assert keys_file == keys_excel


def test_d1_xlsx_returns_ingest_report(
    store, registry, queue, excel_fixture_path, ground_truth
):
    """user story：ingest_file 返回与 ingest_excel 同型的 IngestReport（字段冻结），
    下游台账 / 观测面零改动即可消费。"""
    from ragspine.ingestion.structured.ingestion import ingest_file

    _setup_hk_mapping(store, registry, queue, ground_truth)
    report = ingest_file(excel_fixture_path, store, registry, queue)
    assert isinstance(report, IngestReport)
    assert report.status == "ok"
    assert report.source_doc_id == HK_SCOPE


# ===========================================================================
# D2 — pptx 入库：实体可解析的 pptx 财务表事实写入 fact_store
# ===========================================================================
def test_d2_pptx_facts_ingested(store, registry, queue, tmp_path):
    """user story：作为接手方，我把一张实体可解析（表内 R1C1=ACME Hong Kong）的
    pptx 财务表交给 ingest_file，它应经 pptx_styled_extractor 抽取并把事实写入库，
    让 PPT 来源的数字也能被结构化检索。"""
    from ragspine.ingestion.structured.ingestion import ingest_file

    _init_three(store, registry, queue)
    p = tmp_path / "acme_hk_deck.pptx"
    _make_pptx_financial(p, entity_title=RESOLVABLE_ENTITY_TITLE)
    assert store.count() == 0
    report = ingest_file(p, store, registry, queue)
    assert report.status == "ok"
    assert report.n_facts_ingested >= 1
    assert store.count() >= 1


def test_d2_pptx_fact_queryable(store, registry, queue, tmp_path):
    """user story：pptx 入库后我能按 REVENUE / ACME_HK / FY2024 精确查回该值（2680），
    证明 PPT 事实真正落库且可检索（非仅计数）。"""
    from ragspine.ingestion.structured.ingestion import ingest_file

    _init_three(store, registry, queue)
    p = tmp_path / "acme_hk_deck.pptx"
    _make_pptx_financial(p, entity_title=RESOLVABLE_ENTITY_TITLE)
    ingest_file(p, store, registry, queue)
    rows = store.query("REVENUE", "ACME_HK", "FY", "2024", channel="TOTAL")
    assert len(rows) == 1
    assert rows[0].value == 2680.0


def test_d2_pptx_extractor_version_distinct(store, registry, queue, tmp_path):
    """user story：pptx 入库事实的 extractor_version 应区别于 xlsx（按格式区分血缘），
    审计时一眼可知数字来自 PPT 解析器。"""
    from ragspine.ingestion.structured.ingestion import ingest_file

    _init_three(store, registry, queue)
    p = tmp_path / "acme_hk_deck.pptx"
    _make_pptx_financial(p, entity_title=RESOLVABLE_ENTITY_TITLE)
    ingest_file(p, store, registry, queue)
    rows = store.query("REVENUE", "ACME_HK", "FY", "2024")
    assert len(rows) == 1
    ev = rows[0].extractor_version or ""
    assert "pptx" in ev.lower()
    assert rows[0].review_status in VISIBLE_REVIEW_STATUSES


# ===========================================================================
# D3 — pdf(digital) 入库：pdf_router 判 digital -> 走 pdf_digital_extractor
# ===========================================================================
@pytest.mark.docling  # 唯一真跑 docling 抽取的 dispatch 用例：隔离进程跑（见 test_pdf_digital_extractor）。
def test_d3_pdf_digital_routed_to_digital_extractor(
    store, registry, queue, digital_pdf_path, monkeypatch
):
    """user story：作为接手方，我把数字型 PDF 交给 ingest_file，它应先经 pdf_router
    判定为 digital，再把抽取分发给 pdf_digital_extractor.extract_grids（而非扫描线），
    确保数字型 PDF 走确定性表格解析。"""
    from ragspine.ingestion.structured.ingestion import ingest_file
    import ragspine.ingestion.structured.ingestion as ingestion_mod
    from ragspine.extraction.routing.pdf_router import route, VERDICT_DIGITAL

    # 该 fixture 必须被 router 判为 digital（前置事实）。
    assert route(str(digital_pdf_path)).verdict == VERDICT_DIGITAL

    _init_three(store, registry, queue)

    # spy：分发应调用 pdf_digital_extractor.extract_grids（ingestion 模块内引用）。
    called = {"digital": 0}
    real = None
    if hasattr(ingestion_mod, "pdf_digital_extractor"):
        real = ingestion_mod.pdf_digital_extractor.extract_grids

        def _spy(path):
            called["digital"] += 1
            return real(path)

        monkeypatch.setattr(
            ingestion_mod.pdf_digital_extractor, "extract_grids", _spy
        )

    ingest_file(digital_pdf_path, store, registry, queue)
    assert called["digital"] >= 1


def test_d3_pdf_digital_no_failure(store, registry, queue, digital_pdf_path):
    """user story：数字型 PDF 经 ingest_file 处理后 status 不应为 failed
    （digital 分支正常跑通，不裸抛、不崩）。"""
    from ragspine.ingestion.structured.ingestion import ingest_file

    _init_three(store, registry, queue)
    report = ingest_file(digital_pdf_path, store, registry, queue)
    assert report.status == "ok"
    assert report.error is None
    # 血缘照常记录格式与文件 hash。
    assert report.source_doc_id == "digital.pdf"


# ===========================================================================
# D4 — PDF 越界入队：scanned / ask_for_pptx 不抽数字事实，enqueue 带正确 reason
# ===========================================================================
def test_d4_scanned_pdf_enqueued_not_extracted(
    store, registry, queue, scanned_pdf_path
):
    """user story：作为接手方，我把扫描型 PDF（无文本层）交给 ingest_file，它绝不
    凭空抽数字事实，而应入复核队列（reason 指明需 OCR / 人工复核），让扫描件被
    显式接住而非静默丢弃。"""
    from ragspine.ingestion.structured.ingestion import ingest_file
    from ragspine.extraction.routing.pdf_router import route, VERDICT_SCANNED

    assert route(str(scanned_pdf_path)).verdict == VERDICT_SCANNED

    _init_three(store, registry, queue)
    report = ingest_file(scanned_pdf_path, store, registry, queue)
    # 不抽数字事实
    assert report.n_facts_ingested == 0
    assert store.count() == 0
    # 入队且 reason 指向扫描/OCR/人工复核
    assert report.n_enqueued_review > 0
    pending = queue.list_pending()
    assert len(pending) > 0
    assert any(
        ("扫描" in it.reason) or ("OCR" in it.reason.upper()) or ("复核" in it.reason)
        for it in pending
    )


def test_d4_ppt_export_pdf_asks_for_pptx_and_enqueues(
    store, registry, queue, ppt_export_pdf_path
):
    """user story：PowerPoint 导出的 PDF（producer 命中 PowerPoint）应触发
    ask_for_pptx：ingest_file 不强抽，而是告警 + 入队请求提供 pptx 源
    （PRD「原生优先」），避免在退化 PDF 上做不可靠抽取。"""
    from ragspine.ingestion.structured.ingestion import ingest_file
    from ragspine.extraction.routing.pdf_router import route

    assert route(str(ppt_export_pdf_path)).ask_for_pptx is True

    _init_three(store, registry, queue)
    report = ingest_file(ppt_export_pdf_path, store, registry, queue)
    assert report.n_enqueued_review > 0
    pending = queue.list_pending()
    assert any(
        ("pptx" in it.reason.lower()) or ("PowerPoint" in it.reason) or ("源" in it.reason)
        for it in pending
    )
    # 应至少有一条告警提示原生优先。
    assert any(
        ("pptx" in w.lower()) or ("PowerPoint" in w) for w in report.warnings
    )


def test_d4_scanned_pdf_not_silently_dropped(
    store, registry, queue, scanned_pdf_path
):
    """user story：扫描型 PDF 不得被静默丢弃——即便不入库，也必须留下队列项
    或告警，让运营能看到「这份文件被挡下、原因是什么」。"""
    from ragspine.ingestion.structured.ingestion import ingest_file

    _init_three(store, registry, queue)
    report = ingest_file(scanned_pdf_path, store, registry, queue)
    assert report.n_enqueued_review > 0 or len(report.warnings) > 0


# ===========================================================================
# D5 — 实体不可解析 -> skip + enqueue（关键：不猜实体、不误归因）
# ===========================================================================
def test_d5_unresolvable_entity_skips_ingest(
    store, registry, queue, tmp_path
):
    """user story：作为接手方，当一张表的实体无法归一（A1 与 sheet 名都不是受控实体）
    时，ingest_file 绝不臆造实体（避免 GAP-B 式误归因），该表零事实入库。"""
    from ragspine.ingestion.structured.ingestion import ingest_file

    _init_three(store, registry, queue)
    p = tmp_path / "unknown_entity.xlsx"
    # sheet 名与 A1 都不可归一（'Mystery Co.' 不在 home profile），但指标/期间合法。
    _make_xlsx_financial(p, sheet_title="Sheet1", a1="Mystery Co.")
    report = ingest_file(p, store, registry, queue)
    assert report.n_facts_ingested == 0
    assert store.count() == 0


def test_d5_unresolvable_entity_enqueues_for_human(
    store, registry, queue, tmp_path
):
    """user story：实体不可解析的表必须入复核队列、reason 含「实体无法解析」，
    并带定位血缘（source_doc_id / sheet / locator），让 SME 人工指认而非系统瞎猜。"""
    from ragspine.ingestion.structured.ingestion import ingest_file

    _init_three(store, registry, queue)
    p = tmp_path / "unknown_entity.xlsx"
    _make_xlsx_financial(p, sheet_title="Sheet1", a1="Mystery Co.")
    report = ingest_file(p, store, registry, queue)
    assert report.n_enqueued_review > 0
    pending = queue.list_pending()
    assert any("实体无法解析" in it.reason for it in pending)


# ===========================================================================
# D6 — 无 active 映射 + 有颜色 tag -> enqueue（颜色映射未确认）
# ===========================================================================
def test_d6_colored_no_active_mapping_enqueues(
    store, registry, queue, tmp_path
):
    """user story：作为接手方，当文件确有颜色编码（数值格被着色，需翻译成 tag）但
    该 scope 无 active 颜色映射时，ingest_file 应入复核队列（reason 含「颜色映射未确认」），
    绝不静默把带语义的颜色当无意义忽略。"""
    from ragspine.ingestion.structured.ingestion import ingest_file

    _init_three(store, registry, queue)  # registry 空：无任何 active 映射
    p = tmp_path / "acme_hk_colored.xlsx"
    # 实体可解析（A1='ACME Hong Kong'），B2 数值格上黄色填充（确有颜色 tag 需翻译）。
    _make_xlsx_financial(
        p, sheet_title="HK", a1=RESOLVABLE_ENTITY_TITLE, colored=True
    )
    report = ingest_file(p, store, registry, queue)
    assert report.n_enqueued_review > 0
    pending = queue.list_pending()
    assert any(
        ("颜色映射" in it.reason) or ("映射未确认" in it.reason)
        for it in pending
    )


def test_d6_colored_no_active_mapping_still_warns(
    store, registry, queue, tmp_path
):
    """user story：有色无映射时除入队外仍应在报告里告警（颜色 tag 未翻译），
    确保运营在台账上能看到「这份文件的颜色语义被挂起待确认」。"""
    from ragspine.ingestion.structured.ingestion import ingest_file

    _init_three(store, registry, queue)
    p = tmp_path / "acme_hk_colored.xlsx"
    _make_xlsx_financial(
        p, sheet_title="HK", a1=RESOLVABLE_ENTITY_TITLE, colored=True
    )
    report = ingest_file(p, store, registry, queue)
    assert any(("映射" in w) or ("map" in w.lower()) for w in report.warnings)


# ===========================================================================
# D7 — dry_run：store 与 queue 均零写入，n_facts_ingested / n_enqueued_review 皆 0
# ===========================================================================
def test_d7_dry_run_writes_nothing(
    store, registry, queue, excel_fixture_path, ground_truth
):
    """user story：作为接手方，dry_run 预览只产报告、绝不碰库——fact_store 与
    review_queue 都零写入，n_facts_ingested == 0 且 n_enqueued_review == 0。"""
    from ragspine.ingestion.structured.ingestion import ingest_file

    _setup_hk_mapping(store, registry, queue, ground_truth)
    report = ingest_file(
        excel_fixture_path, store, registry, queue, dry_run=True
    )
    assert report.dry_run is True
    assert report.n_facts_ingested == 0
    assert report.n_enqueued_review == 0
    assert store.count() == 0
    assert queue.list_pending() == []


def test_d7_dry_run_report_otherwise_complete(
    store, registry, queue, excel_fixture_path, ground_truth
):
    """user story：dry_run 仍完整跑抽取——报告其余字段（n_grids / n_facts_extracted /
    file_hash）照常产出，让我能在不写库的前提下预览将会入库的内容。"""
    from ragspine.ingestion.structured.ingestion import ingest_file

    _setup_hk_mapping(store, registry, queue, ground_truth)
    report = ingest_file(
        excel_fixture_path, store, registry, queue, dry_run=True
    )
    assert report.n_grids >= 1
    assert report.n_facts_extracted >= EXPECTED_HK_FACTS
    assert report.file_hash


def test_d7_dry_run_scanned_pdf_no_queue_write(
    store, registry, queue, scanned_pdf_path
):
    """user story：扫描型 PDF 的 dry_run 也不得写队列——dry_run 是「只看不动」的硬约定，
    即便走的是入队分支，n_enqueued_review 仍为 0、队列保持空。"""
    from ragspine.ingestion.structured.ingestion import ingest_file

    _init_three(store, registry, queue)
    report = ingest_file(
        scanned_pdf_path, store, registry, queue, dry_run=True
    )
    assert report.n_enqueued_review == 0
    assert queue.list_pending() == []


# ===========================================================================
# D8 — 幂等：同文件 ingest_file 两次，store.count() 不增长
# ===========================================================================
def test_d8_idempotent_xlsx(
    store, registry, queue, excel_fixture_path, ground_truth
):
    """user story：作为接手方，同一文件重复 ingest_file 必须幂等——靠 fact_metric
    唯一键 upsert，库内事实总数不随重复入库增长。"""
    from ragspine.ingestion.structured.ingestion import ingest_file

    _setup_hk_mapping(store, registry, queue, ground_truth)
    ingest_file(excel_fixture_path, store, registry, queue)
    first = store.count()
    ingest_file(excel_fixture_path, store, registry, queue)
    assert store.count() == first


def test_d8_idempotent_pptx(store, registry, queue, tmp_path):
    """user story：pptx 同样满足幂等——重复 ingest_file 同一 pptx 财务表，
    同一指标×实体×期间×渠道仍只有 1 条。"""
    from ragspine.ingestion.structured.ingestion import ingest_file

    _init_three(store, registry, queue)
    p = tmp_path / "acme_hk_deck.pptx"
    _make_pptx_financial(p, entity_title=RESOLVABLE_ENTITY_TITLE)
    ingest_file(p, store, registry, queue)
    first = store.count()
    ingest_file(p, store, registry, queue)
    assert store.count() == first
    rows = store.query("REVENUE", "ACME_HK", "FY", "2024")
    assert len(rows) == 1


# ===========================================================================
# D9 — 不支持后缀：status='failed' / 妥善处理，不崩
# ===========================================================================
def test_d9_unsupported_suffix_does_not_raise(
    store, registry, queue, tmp_path
):
    """user story：作为接手方，把不支持的格式（.txt）交给 ingest_file，它绝不裸抛
    异常把整批拖垮，而应返回一份报告（status='failed' 并记 error）。"""
    from ragspine.ingestion.structured.ingestion import ingest_file

    _init_three(store, registry, queue)
    p = tmp_path / "notes.txt"
    p.write_text("这不是一个可抽取的表格文件", encoding="utf-8")
    report = ingest_file(p, store, registry, queue)  # 不应抛异常
    assert isinstance(report, IngestReport)
    assert report.status == "failed"
    assert report.error


def test_d9_unsupported_suffix_writes_nothing(
    store, registry, queue, tmp_path
):
    """user story：不支持后缀的文件不得污染库——fact_store 零写入。"""
    from ragspine.ingestion.structured.ingestion import ingest_file

    _init_three(store, registry, queue)
    p = tmp_path / "notes.txt"
    p.write_text("plain text, no table", encoding="utf-8")
    ingest_file(p, store, registry, queue)
    assert store.count() == 0


# ===========================================================================
# D10 — 回归：ingest_excel 既有行为不变
# ===========================================================================
def test_d10_ingest_excel_unchanged_fact_count(
    store, registry, queue, excel_fixture_path, ground_truth
):
    """user story：接通分发入口后，既有 ingest_excel 行为零回归——对 fixture 仍入
    12 条事实、status='ok'。"""
    _setup_hk_mapping(store, registry, queue, ground_truth)
    report = ingest_excel(excel_fixture_path, store, registry, queue)
    assert report.status == "ok"
    assert report.n_facts_ingested == EXPECTED_HK_FACTS
    assert store.count() == EXPECTED_HK_FACTS


def test_d10_ingest_excel_unchanged_no_spurious_enqueue(
    store, registry, queue, excel_fixture_path, ground_truth
):
    """user story：既有 xlsx fixture（HK_Performance 实体可解析、有 active 映射）
    仍不触发任何新增 enqueue 条件——n_enqueued_review 保持 0，证明接线没把能确定
    的情况误发去复核。"""
    _setup_hk_mapping(store, registry, queue, ground_truth)
    report = ingest_excel(excel_fixture_path, store, registry, queue)
    assert report.n_enqueued_review == 0
    assert queue.list_pending() == []
