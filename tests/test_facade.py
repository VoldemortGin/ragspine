"""High-level :class:`RAGSpine` facade contracts.

These tests deliberately exercise only the public facade.  The low-level ingestion,
storage, and agent tests remain responsible for their individual domain rules.
"""

from pathlib import Path

import pytest
from openpyxl import Workbook
from pptx import Presentation
from pptx.util import Inches


def _make_financial_workbook(path: Path) -> None:
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "ACME Hong Kong"
    sheet["A1"] = "ACME Hong Kong"
    sheet["B1"] = "FY2024"
    sheet["A2"] = "REVENUE"
    sheet["B2"] = 2680
    workbook.save(path)


def _make_dual_channel_deck(path: Path) -> None:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    table = slide.shapes.add_table(2, 2, Inches(0.5), Inches(0.5), Inches(8), Inches(2)).table
    table.cell(0, 0).text = "ACME Hong Kong"
    table.cell(0, 1).text = "FY2024"
    table.cell(1, 0).text = "REVENUE"
    table.cell(1, 1).text = "2680"
    textbox = slide.shapes.add_textbox(Inches(0.5), Inches(3), Inches(8), Inches(1))
    textbox.text_frame.text = "Revenue improved after the agency channel expansion."
    presentation.save(path)


def test_local_context_owns_and_closes_workspace_resources(tmp_path):
    """A local workspace is usable in its context and rejects work after closing."""
    from ragspine import RAGSpine

    workspace = tmp_path / "knowledge-base"
    with RAGSpine.local(workspace) as rag:
        assert workspace.is_dir()
        assert (workspace / "knowledge.db").is_file()
        assert (workspace / "mapping.db").is_file()
        assert (workspace / "review.db").is_file()

    with pytest.raises(RuntimeError, match="closed|关闭"):
        rag.ask("香港 FY2024 REVENUE 是多少")


def test_local_defaults_to_lean_economy_profile(tmp_path):
    """The zero-configuration facade never opts into model-backed retrieval."""
    from ragspine import RAGSpine
    from ragspine.facade import RetrievalProfile, make_retrieval_preset

    with RAGSpine.local(tmp_path / "knowledge-base") as rag:
        assert rag.retrieval == make_retrieval_preset(RetrievalProfile.ECONOMY)


def test_balanced_profile_drives_offline_narrative_ask(tmp_path):
    """Balanced retrieval is assembled by ask and remains deterministic offline."""
    from ragspine import RAGSpine
    from ragspine.facade import RetrievalProfile

    source = tmp_path / "review.txt"
    source.write_text("营收增长的原因是代理渠道扩张。", encoding="utf-8")

    with RAGSpine.local(tmp_path / "knowledge-base", profile=RetrievalProfile.BALANCED) as rag:
        rag.ingest(source)
        result = rag.ask("营收为什么增长？")

    assert result.sources
    assert result.sources[0]["doc"] == "review.txt"


def test_balanced_preset_exposes_the_plan_used_by_the_facade(tmp_path):
    """The public preset spelling resolves once and drives runtime assembly."""
    from ragspine import RAGSpine

    with RAGSpine.local(tmp_path / "knowledge-base", preset="balanced") as rag:
        assert rag.retrieval.retrieval_mode == "hybrid"
        assert rag.retrieval.embedding == "deterministic"
        assert rag.retrieval.vector_store == "in_process"
        assert rag.effective_plan.config.profile == "balanced"
        assert rag.effective_plan.config.retrieval.embedding == "deterministic"
        assert rag.effective_plan.source_for("profile") == "preset"


def test_preset_rejects_the_legacy_profile_spelling(tmp_path):
    """Two names for the same selection are rejected instead of guessed."""
    from ragspine import RAGSpine

    with pytest.raises(ValueError, match="preset.*profile"):
        RAGSpine.local(tmp_path / "knowledge-base", preset="balanced", profile="quality")


def test_explicit_retrieval_preset_overrides_named_profile(tmp_path):
    """Advanced callers can replace a profile without untyped configuration maps."""
    from ragspine import RAGSpine
    from ragspine.facade import RetrievalProfile, make_retrieval_preset

    custom = make_retrieval_preset(
        RetrievalProfile.BALANCED,
        embedding="none",
        vector_store="none",
    )
    with RAGSpine.local(
        tmp_path / "knowledge-base",
        profile=RetrievalProfile.QUALITY,
        retrieval=custom,
    ) as rag:
        assert rag.retrieval is custom


def test_ingest_then_ask_preserves_answer_provenance_across_reopen(tmp_path):
    """User-owned structured data remains answerable after reopening a workspace."""
    from ragspine import RAGSpine

    workspace = tmp_path / "knowledge-base"
    source = tmp_path / "acme_hk_fy2024.xlsx"
    _make_financial_workbook(source)

    with RAGSpine.local(workspace) as rag:
        ingest_result = rag.ingest(source)
        assert not ingest_result.failed

    with RAGSpine.local(workspace) as rag:
        answer = rag.ask("香港 FY2024 REVENUE 是多少")

    assert "2680" in answer.answer
    assert answer.sources == [
        {
            "doc": "acme_hk_fy2024.xlsx",
            "locator": "sheet=ACME Hong Kong!B2",
        }
    ]


def test_xlsx_ingest_returns_one_structured_aggregate_result(tmp_path):
    """XLSX defaults to the structured channel behind one stable result type."""
    from ragspine import RAGSpine
    from ragspine.facade import IngestResult

    source = tmp_path / "financials.xlsx"
    _make_financial_workbook(source)

    with RAGSpine.local(tmp_path / "knowledge-base") as rag:
        result = rag.ingest(source)

    assert isinstance(result, IngestResult)
    assert len(result.structured_reports) == 1
    assert result.structured_reports[0].status == "ok"
    assert result.narrative_report is None
    assert not result.failed
    assert isinstance(result.summary, str) and result.summary


def test_pptx_ingest_attempts_both_channels_without_cross_channel_failure(tmp_path):
    """PPTX keeps structured facts and narrative chunks through one dispatch."""
    from ragspine import RAGSpine
    from ragspine.facade import IngestResult

    source = tmp_path / "board_pack_FY2024.pptx"
    _make_dual_channel_deck(source)

    with RAGSpine.local(tmp_path / "knowledge-base") as rag:
        result = rag.ingest(source)

    assert isinstance(result, IngestResult)
    assert len(result.structured_reports) == 1
    assert result.structured_reports[0].status == "ok"
    assert result.structured_reports[0].n_facts_ingested >= 1
    assert result.narrative_report is not None
    assert result.narrative_report.files[0].n_chunks >= 1
    assert not result.failed


def test_ingest_routes_text_to_narrative_channel(tmp_path):
    """Plain text takes the narrative path without requiring structured extraction."""
    from ragspine import RAGSpine

    source = tmp_path / "notes.txt"
    source.write_text("Revenue narrative", encoding="utf-8")

    with RAGSpine.local(tmp_path / "knowledge-base") as rag:
        result = rag.ingest(source)

    assert result.structured_reports == ()
    assert result.narrative_report is not None
    assert result.narrative_report.files[0].n_chunks == 1
    assert not result.failed


def test_parent_child_config_expands_a_child_hit_to_its_parent_context(tmp_path):
    """The high-level indexing choice changes real ingest and generation behavior."""
    from ragspine import RAGSpine

    source = tmp_path / "review.txt"
    source.write_text(
        "# 营收\n营收增长来自渠道扩张。\n毛利改善来自产品组合。\n# 风险\n汇率仍有波动。",
        encoding="utf-8",
    )
    config = {
        "indexing": {"chunker": "parent_child", "max_chars": 16, "overlap_chars": 0}
    }

    with RAGSpine.local(tmp_path / "knowledge-base", config=config) as rag:
        rag.ingest(source)
        result = rag.ask("营收为什么增长？")

    assert "营收增长来自渠道扩张" in result.answer_plain
    assert "毛利改善来自产品组合" in result.answer_plain
    assert result.sources
    assert result.sources[0]["doc"] == "review.txt"
    assert any(source["locator"].endswith("#para1-2") for source in result.sources)


def test_reopen_refuses_to_query_with_an_incompatible_indexing_contract(tmp_path):
    """A workspace never silently reuses chunks built with another contract."""
    from ragspine import RAGSpine
    from ragspine.config import ReindexRequiredError

    workspace = tmp_path / "knowledge-base"
    source = tmp_path / "review.txt"
    source.write_text("# 营收\n营收增长来自渠道扩张。\n毛利改善。", encoding="utf-8")
    parent_child = {
        "indexing": {"chunker": "parent_child", "max_chars": 16, "overlap_chars": 0}
    }

    with RAGSpine.local(workspace, config=parent_child) as rag:
        rag.ingest(source)

    with RAGSpine.local(workspace) as rag:
        with pytest.raises(ReindexRequiredError, match=r"chunking.*ragspine ingest.*--reindex"):
            rag.ask("营收为什么增长？")


def test_reopen_refuses_incremental_ingest_with_an_incompatible_contract(tmp_path):
    """Compatibility is checked before an incremental writer can mutate the index."""
    from ragspine import RAGSpine
    from ragspine.config import ReindexRequiredError

    workspace = tmp_path / "knowledge-base"
    first = tmp_path / "first.txt"
    second = tmp_path / "second.txt"
    first.write_text("First narrative.", encoding="utf-8")
    second.write_text("Second narrative.", encoding="utf-8")
    parent_child = {"indexing": {"chunker": "parent_child"}}
    with RAGSpine.local(workspace, config=parent_child) as rag:
        rag.ingest(first)

    with RAGSpine.local(workspace) as incompatible:
        with pytest.raises(ReindexRequiredError, match="chunking"):
            incompatible.ingest(second)

    with RAGSpine.local(workspace, config=parent_child) as compatible:
        assert compatible.ask("Why did the first narrative change?").sources[0]["doc"] == "first.txt"


def test_structured_only_ingest_ignores_narrative_index_contract_mismatch(tmp_path):
    """A narrative fingerprint never blocks an independent structured writer."""
    from ragspine import RAGSpine

    workspace = tmp_path / "knowledge-base"
    narrative = tmp_path / "review.txt"
    workbook = tmp_path / "financials.xlsx"
    narrative.write_text("Revenue improved after channel expansion.", encoding="utf-8")
    _make_financial_workbook(workbook)

    with RAGSpine.local(workspace) as rag:
        rag.ingest(narrative)

    with RAGSpine.local(
        workspace, config={"indexing": {"chunker": "parent_child"}}
    ) as incompatible_for_narrative:
        result = incompatible_for_narrative.ingest(workbook)

    assert result.narrative_report is None
    assert len(result.structured_reports) == 1
    assert result.structured_reports[0].status == "ok"
    assert result.structured_reports[0].n_facts_ingested >= 1


def test_parent_child_context_survives_a_compatible_reopen(tmp_path):
    """Reopening reconstructs parent context from the persisted public workspace."""
    from ragspine import RAGSpine

    workspace = tmp_path / "knowledge-base"
    source = tmp_path / "review.txt"
    source.write_text(
        "# 营收\n营收增长来自渠道扩张。\n毛利改善来自产品组合。", encoding="utf-8"
    )
    config = {
        "indexing": {"chunker": "parent_child", "max_chars": 16, "overlap_chars": 0}
    }
    with RAGSpine.local(workspace, config=config) as rag:
        rag.ingest(source)

    with RAGSpine.local(workspace, config=config) as reopened:
        result = reopened.ask("营收为什么增长？")

    assert "毛利改善来自产品组合" in result.answer_plain
    assert any(source["locator"].endswith("#para1-2") for source in result.sources)


def test_dry_run_does_not_claim_the_workspace_index_contract(tmp_path):
    """A preview leaves the workspace available for a different real contract."""
    from ragspine import RAGSpine

    workspace = tmp_path / "knowledge-base"
    source = tmp_path / "review.txt"
    source.write_text("Revenue improved after channel expansion.", encoding="utf-8")
    parent_child = {
        "indexing": {"chunker": "parent_child", "max_chars": 16, "overlap_chars": 0}
    }

    with RAGSpine.local(workspace, config=parent_child) as preview:
        report = preview.ingest(source, dry_run=True)
        assert report.narrative_report is not None
        assert report.narrative_report.dry_run is True
        assert report.narrative_report.files[0].n_chunks > 0

    with RAGSpine.local(workspace) as real:
        report = real.ingest(source)
        result = real.ask("Why did revenue improve?")

    assert report.narrative_report is not None
    assert report.narrative_report.files[0].status == "ingested"
    assert result.sources[0]["doc"] == "review.txt"


def test_explicit_default_indexing_is_equivalent_to_zero_configuration(tmp_path):
    """Naming the canonical default does not change public ingest or ask results."""
    from ragspine import RAGSpine

    source = tmp_path / "review.txt"
    source.write_text("Revenue improved after channel expansion.", encoding="utf-8")
    with RAGSpine.local(tmp_path / "implicit") as implicit:
        implicit_report = implicit.ingest(source)
        implicit_result = implicit.ask("Why did revenue improve?")
    with RAGSpine.local(
        tmp_path / "explicit",
        config={"indexing": {"chunker": "none", "max_chars": 480, "overlap_chars": 80}},
    ) as explicit:
        explicit_report = explicit.ingest(source)
        explicit_result = explicit.ask("Why did revenue improve?")

    assert implicit_report.summary == explicit_report.summary
    assert implicit_result.answer_plain == explicit_result.answer_plain
    assert implicit_result.sources == explicit_result.sources
