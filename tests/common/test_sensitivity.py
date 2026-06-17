"""叙事入库敏感度【确定性自动分级】单测（TDD 红色阶段）。

安全 P0：现状 narrative_ingest 对未显式标 sensitivity 的文档一律落 'INTERNAL'，
RESTRICTED 文档若漏标即以 INTERNAL 流进检索/prompt/回答，两道出域过滤
（narrative_link / listwise_rerank 剔除 RESTRICTED）形同虚设——漏标 = 泄露。

本测试钉死的契约：
    A) 敏感规则配置化：config 的 [sensitivity] 段 → SensitivityPolicy（frozen），
       由 load_company_profile 挂到 CompanyProfile.sensitivity；缺省回退内置默认
       （default_level='INTERNAL'，保证既有行为不变）。
    B) 分级器 ragspine.common.sensitivity.classify_sensitivity(filename, text, policy)：
       文件名/路径命中受限模式 → RESTRICTED；否则正文命中受限关键词 → RESTRICTED；
       否则 escalate_unknown_to_restricted 为 True 时 → RESTRICTED，否则 default_level。
       全大小写不敏感、确定性、零外部调用。
    C) 接线 narrative_ingest：meta 未显式给 sensitivity 时调分级器决定；
       显式给则尊重显式值（人工标注优先）。
    D) 端到端泄露守护：带 RESTRICTED 信号、未显式标的 doc 过 ingest →
       块 sensitivity=='RESTRICTED' → 经出域过滤后不出现在检索结果/喂 prompt 候选里。

红色预期：ragspine.common.sensitivity 尚不存在、CompanyProfile 无 sensitivity 字段、
narrative_ingest 未接线分级器 —— 这些导入/断言失败即为「正确原因的红」。
"""

import os

import pytest
import rootutils

ROOT_DIR = rootutils.setup_root(os.getcwd(), indicator=".project-root", pythonpath=True)

from ragspine.retrieval.chunking.chunk_store import ChunkStore
from ragspine.common.company_profile import CompanyProfile, load_company_profile
from ragspine.common.sensitivity import SensitivityPolicy, classify_sensitivity


# ---------------------------------------------------------------------------
# 测试用 policy 构造（不依赖部署配置文件，规则就地声明，便于断言 config 驱动）
# ---------------------------------------------------------------------------

def _acme_policy(*, escalate: bool = False) -> SensitivityPolicy:
    """模拟 ACME 部署的 [sensitivity]：中性占位受限文件名/关键词。

    文件名模式覆盖 exco / pr_rating / board_minutes；关键词覆盖
    高管绩效评级 / remuneration / exec_comp。刻意不含 '董事会'，以反证
    分级器完全由传入 policy 决定（见 test_s6_acme_only_keyword_not_restricted）。
    """
    return SensitivityPolicy(
        default_level="INTERNAL",
        escalate_unknown_to_restricted=escalate,
        restricted_filename_patterns=["exco", "pr_rating", "board_minutes"],
        restricted_keywords=["高管绩效评级", "remuneration", "exec_comp"],
    )


# ===========================================================================
# S1 文件名信号
# ===========================================================================

def test_s1_filename_signal_escalates_to_restricted():
    """user story：作为安全责任人，一份文件名含 'Exco' 的纪要即便正文普通、
    入库时漏标 sensitivity，分级器也必须凭文件名判定 RESTRICTED，避免漏标泄露。"""
    level = classify_sensitivity(
        "Exco_Minutes_2025Q4.pptx",
        "本季度业务进展顺利，渠道扩张按计划推进。",
        _acme_policy(),
    )
    assert level == "RESTRICTED"


def test_s1_filename_pattern_is_case_insensitive_substring():
    """user story：受限文件名模式按大小写不敏感子串匹配，
    'PR_RATING' 与配置里的 'pr_rating' 等价命中，不因大小写漏判。"""
    level = classify_sensitivity(
        "GCE_PR_RATING_FY2025.pptx", "普通正文", _acme_policy()
    )
    assert level == "RESTRICTED"


# ===========================================================================
# S2 正文关键词信号
# ===========================================================================

def test_s2_body_keyword_signal_escalates_to_restricted():
    """user story：文件名看似普通，但正文出现'高管绩效评级'这类受限关键词时，
    分级器必须升级为 RESTRICTED，堵住"换个文件名就泄露"的口子。"""
    level = classify_sensitivity(
        "MBR_Summary.pptx",
        "本页讨论高管绩效评级与下一阶段安排。",
        _acme_policy(),
    )
    assert level == "RESTRICTED"


def test_s2_body_keyword_is_case_insensitive():
    """user story：英文受限关键词大小写不敏感命中（'REMUNERATION' ≡ 'remuneration'）。"""
    level = classify_sensitivity(
        "Notes.pdf",
        "The REMUNERATION committee reviewed the package.",
        _acme_policy(),
    )
    assert level == "RESTRICTED"


# ===========================================================================
# S3 不过度分级（关键回归）
# ===========================================================================

def test_s3_normal_financial_report_stays_internal():
    """user story：作为检索质量负责人，普通财报（文件名/正文都无受限信号）
    绝不能被误升级为 RESTRICTED，否则会被出域过滤藏掉、击穿 41 golden 与检索。"""
    level = classify_sensitivity(
        "ACME_FY2025_Results.pptx",
        "集团 REVENUE 同比增长，香港与中国内地为主要驱动。",
        _acme_policy(),
    )
    assert level == "INTERNAL"


# ===========================================================================
# S5 strict 开关
# ===========================================================================

def test_s5_strict_switch_escalates_unsignaled_docs():
    """user story：在 strict 部署下（escalate_unknown_to_restricted=True），
    任何无明确信号的文档都 fail-safe 到 RESTRICTED；本开关默认 false 以保行为。"""
    no_signal = classify_sensitivity(
        "Generic_Update.pptx", "一般业务更新，无敏感内容。", _acme_policy(escalate=True)
    )
    assert no_signal == "RESTRICTED"


def test_s5_strict_switch_default_false_keeps_internal():
    """user story：strict 开关默认关闭——无信号文档仍走 default_level（INTERNAL），
    证明 blanket 升级只是可选开关、不是默认行为。"""
    assert _acme_policy().escalate_unknown_to_restricted is False
    no_signal = classify_sensitivity(
        "Generic_Update.pptx", "一般业务更新，无敏感内容。", _acme_policy()
    )
    assert no_signal == "INTERNAL"


# ===========================================================================
# S6 配置/泛化：规则真由 config 驱动，非硬编码
# ===========================================================================

def test_s6_acme_specific_keyword_escalates_under_acme():
    """user story：ACME 部署用自己的受限词（'exec_comp'），命中即 RESTRICTED——
    证明规则来自 config 而非任何写死的 ACME 词表。"""
    level = classify_sensitivity(
        "Q3_Pack.pptx", "see exec_comp appendix for detail", _acme_policy()
    )
    assert level == "RESTRICTED"


def test_s6_acme_only_keyword_not_restricted_under_acme():
    """user story：仅在 ACME 配置里受限的词（如 '董事会'）在 ACME 配置下不属受限词，
    不应升级 —— 反证分级器没有硬编码 ACME 词、完全由传入 policy 决定。"""
    level = classify_sensitivity(
        "Q3_Pack.pptx", "董事会上讨论了常规事项。", _acme_policy()
    )
    assert level == "INTERNAL"


def test_s6_default_level_respected_per_policy():
    """user story：default_level 由 policy 决定，无信号文档落 policy.default_level。"""
    pub_policy = SensitivityPolicy(
        default_level="PUBLIC",
        escalate_unknown_to_restricted=False,
        restricted_filename_patterns=[],
        restricted_keywords=[],
    )
    assert classify_sensitivity("Press_Release.pdf", "公开新闻稿。", pub_policy) == "PUBLIC"


# ===========================================================================
# A) 配置承载：CompanyProfile.sensitivity + load_company_profile
# ===========================================================================

def test_company_profile_carries_sensitivity_policy():
    """user story：部署 config/company.toml 的 [sensitivity] 段被解析进
    CompanyProfile.sensitivity（SensitivityPolicy 实例）。"""
    profile = load_company_profile()
    assert isinstance(profile, CompanyProfile)
    assert isinstance(profile.sensitivity, SensitivityPolicy)
    assert profile.sensitivity.default_level == "INTERNAL"


def test_sensitivity_policy_is_frozen():
    """user story：SensitivityPolicy 不可变（frozen dataclass），防运行期被改。"""
    policy = _acme_policy()
    with pytest.raises(Exception):
        policy.default_level = "MUTATED"  # type: ignore[misc]


def test_missing_file_falls_back_to_builtin_sensitivity_default(tmp_path):
    """user story：配置文件缺失时静默回退内置默认敏感策略，
    default_level='INTERNAL'、strict 开关 False —— 保证既有行为字节级不变。"""
    profile = load_company_profile(tmp_path / "no_such.toml")
    assert profile.sensitivity.default_level == "INTERNAL"
    assert profile.sensitivity.escalate_unknown_to_restricted is False


def test_sensitivity_policy_loaded_from_toml(tmp_path):
    """user story：临时 toml 的 [sensitivity] 段被完整读入 policy 字段，
    证明四个字段均由 config 驱动。"""
    toml_text = (
        "[home]\n"
        'company_name = "ACME"\n'
        'entity_code = "ACME_GROUP"\n'
        "\n"
        "[sensitivity]\n"
        'default_level = "INTERNAL"\n'
        "escalate_unknown_to_restricted = true\n"
        'restricted_filename_patterns = ["board_minutes", "exec_comp"]\n'
        'restricted_keywords = ["绩效", "exec_comp"]\n'
    )
    path = tmp_path / "company.toml"
    path.write_text(toml_text, encoding="utf-8")

    profile = load_company_profile(path)
    pol = profile.sensitivity
    assert pol.default_level == "INTERNAL"
    assert pol.escalate_unknown_to_restricted is True
    assert "board_minutes" in pol.restricted_filename_patterns
    assert "exec_comp" in pol.restricted_keywords
    # 端到端：该 ACME policy 驱动分级器
    assert classify_sensitivity("Board_Minutes_Q1.pdf", "x", pol) == "RESTRICTED"


# ===========================================================================
# C) 接线 narrative_ingest：分级器决定 / 显式标注优先
# ===========================================================================

def _make_deck(path, paragraphs: list[str], notes: str | None = None) -> None:
    """单页 deck：每段一个文本框（与 test_narrative_ingest 同构造）。"""
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    for i, text in enumerate(paragraphs):
        tb = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.5 + i * 0.9), Inches(8), Inches(0.8)
        )
        tb.text_frame.text = text
    if notes:
        slide.notes_slide.notes_text_frame.text = notes
    prs.save(str(path))


def _store(tmp_path):
    store = ChunkStore(tmp_path / "chunks.db")
    store.init_schema()
    return store


def _doc_sensitivity(store: ChunkStore, doc_id: str) -> set[str]:
    """取某 doc 活跃块的 sensitivity 取值集合。"""
    chunks = store.iter_chunks(doc_id=doc_id)
    return {c.sensitivity for c in chunks}


def test_s4_explicit_meta_sensitivity_wins(tmp_path, monkeypatch):
    """user story：人工显式标注优先于分级器——meta 给 sensitivity='PUBLIC' 时，
    即便文件名/正文有受限信号，入库结果也尊重 PUBLIC（人工判断高于规则）。"""
    import ragspine.ingestion.narrative.narrative_ingest as ni

    # 用带 ACME 受限规则的 profile，确保若分级器越权会被它升级（从而暴露 bug）
    monkeypatch.setattr(
        ni, "_PROFILE",
        _profile_with_policy(_acme_policy()),
        raising=False,
    )
    deck = tmp_path / "Exco_Notes.pptx"
    _make_deck(deck, ["本页讨论高管绩效评级。"])
    store = _store(tmp_path)
    try:
        ni.ingest_narrative(
            [deck], store, meta_by_doc={"Exco_Notes.pptx": {"sensitivity": "PUBLIC"}}
        )
        assert _doc_sensitivity(store, "Exco_Notes.pptx") == {"PUBLIC"}
    finally:
        store.close()


def test_c_unlabeled_restricted_doc_is_auto_classified(tmp_path, monkeypatch):
    """user story：未显式标 sensitivity 的受限文档（文件名含 Exco），
    入库时分级器自动判 RESTRICTED 并写入块——堵住漏标=泄露。"""
    import ragspine.ingestion.narrative.narrative_ingest as ni

    monkeypatch.setattr(
        ni, "_PROFILE", _profile_with_policy(_acme_policy()), raising=False
    )
    deck = tmp_path / "Exco_Minutes.pptx"
    _make_deck(deck, ["香港业务进展讨论。"])
    store = _store(tmp_path)
    try:
        ni.ingest_narrative([deck], store)
        assert _doc_sensitivity(store, "Exco_Minutes.pptx") == {"RESTRICTED"}
    finally:
        store.close()


def test_s8_normal_internal_doc_unchanged(tmp_path, monkeypatch):
    """user story（回归）：普通 INTERNAL 文档（无信号、未显式标）入库后
    sensitivity 仍为 'INTERNAL'，与既有行为一致，不被误升级。"""
    import ragspine.ingestion.narrative.narrative_ingest as ni

    monkeypatch.setattr(
        ni, "_PROFILE", _profile_with_policy(_acme_policy()), raising=False
    )
    deck = tmp_path / "ACME_FY2025_Results.pptx"
    _make_deck(deck, ["集团 REVENUE 同比增长，香港与中国内地为主要驱动。"])
    store = _store(tmp_path)
    try:
        ni.ingest_narrative([deck], store)
        assert _doc_sensitivity(store, "ACME_FY2025_Results.pptx") == {"INTERNAL"}
    finally:
        store.close()


# ===========================================================================
# S7 端到端泄露守护（最重要）
# ===========================================================================

def test_s7_unlabeled_restricted_doc_never_egresses(tmp_path, monkeypatch):
    """user story：作为安全责任人，一份带 RESTRICTED 信号但漏标 sensitivity 的
    文档过 narrative_ingest 后，其块必须被自动判为 RESTRICTED，并经
    NarrativeIndexRetriever 出域过滤后【绝不出现】在喂 LLM 的候选/最终来源里。"""
    import ragspine.ingestion.narrative.narrative_ingest as ni
    from ragspine.retrieval.link.narrative_link import NarrativeIndexRetriever
    from ragspine.retrieval.lexical.retrieval import NarrativeIndex

    monkeypatch.setattr(
        ni, "_PROFILE", _profile_with_policy(_acme_policy()), raising=False
    )
    # 一份普通文档（应被检索到）+ 一份漏标的受限文档（含独特泄露标记）
    normal = tmp_path / "HK_QBR_2025Q4.pptx"
    _make_deck(normal, ["香港 REVENUE 下降主因是 MCV 客群收缩与银保渠道调整。"])
    leaky = tmp_path / "Exco_Minutes.pptx"
    _make_deck(leaky, ["香港 REVENUE 下降的高管绩效评级讨论 SECRET_TOKEN。"])

    store = _store(tmp_path)
    try:
        ni.ingest_narrative([normal, leaky], store)
        # 漏标文档已被自动判 RESTRICTED
        assert _doc_sensitivity(store, "Exco_Minutes.pptx") == {"RESTRICTED"}

        index = NarrativeIndex(store)
        adapter = NarrativeIndexRetriever(index)
        snippets = adapter.retrieve("香港REVENUE为什么下降")
        assert snippets  # 普通文档应被召回
        assert all(s["doc_id"] != "Exco_Minutes.pptx" for s in snippets)
        assert all("SECRET_TOKEN" not in s["text"] for s in snippets)
    finally:
        store.close()


# ---------------------------------------------------------------------------
# 辅助：把一个 SensitivityPolicy 包成带该策略的 CompanyProfile（其余字段取默认）
# ---------------------------------------------------------------------------

def _profile_with_policy(policy: SensitivityPolicy) -> CompanyProfile:
    base = load_company_profile("definitely-not-a-real-path.toml")
    return CompanyProfile(
        home_company_name=base.home_company_name,
        home_entity_code=base.home_entity_code,
        home_entity_synonyms=dict(base.home_entity_synonyms),
        entity_geography=dict(base.entity_geography),
        external_entities=dict(base.external_entities),
        home_entity_labels=dict(base.home_entity_labels),
        sensitivity=policy,
    )
