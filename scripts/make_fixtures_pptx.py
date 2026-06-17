"""生成三期（PPT 增强 + 扫描线）合成 fixture + ground truth（确定性、硬编码）。

覆盖 PRD 三期刁钻形态（user stories 10、12、13、14）：
    slide1: 标题 'ACME Hong Kong — Product Mix' + 带填充色的表格（行=产品线，
            列=FY2022–2024）。填充色编码属性：
                · 部分行黄色 FFFF00 = 新产品线；
                · 部分行绿色 92D050 = 成熟产品线；
                · 至少一格 theme 色（accent1，经 ppt/theme1.xml 解析）；
                · 至少一格无填充（resolved_rgb=None）。
            外加一个含数字的文本框（'FY2024 REVENUE reached US$2,680m, up 14% YoY'）。
    slide2: 转置表（期间在行、指标在列）—— 不规则形态（语义判断归下游）。
    slide2 演讲者备注：'PROFIT for FY2024 was US$6,950m pending final audit'。

产出：
    data/fixtures/pptx/styled_deck.pptx
    data/fixtures/pptx/pptx_ground_truth.json
        —— 逐格值 / 填充 RGB / 表维度；note fragments 期望（text/source_kind/
           glossary_hits）；以及【OCR fake 测试向量】：针对既有
           data/fixtures/pdf/scanned.pdf 的 3 页，给出 FakeBackend 应返回的
           OcrPageResult 数据（含一张 4×4 表、混入 2 个低置信格 confidence<0.85）
           与期望的网格 / 入队结果。scanned.pdf 的页面本就是 digital.pdf 表格页的
           渲染图，真值沿用其表格数值（REVENUE/NEWSALES/PROFIT × FY2022–2024）。

末尾用 python-pptx 读回自校验（值 / 填充 / 备注写入成功）。
从项目根目录运行：.venv/bin/python scripts/make_fixtures_pptx.py
"""

import json
import os

import rootutils

ROOT_DIR = rootutils.setup_root(os.getcwd(), indicator=".project-root", pythonpath=True)

import xml.etree.ElementTree as ET

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_COLOR_TYPE, MSO_THEME_COLOR
from pptx.util import Inches

OUT_DIR = ROOT_DIR / "data" / "fixtures" / "pptx"
PPTX_PATH = OUT_DIR / "styled_deck.pptx"
GT_PATH = OUT_DIR / "pptx_ground_truth.json"

# scanned.pdf（二期已生成的扫描型 fixture）路径，OCR fake 向量引用它。
SCANNED_PDF_REL = "data/fixtures/pdf/scanned.pdf"
SCANNED_PDF_PATH = ROOT_DIR / "data" / "fixtures" / "pdf" / "scanned.pdf"

# --- 语义色（与 Excel 线同一套：颜色编码属性，user story 13）-----------------
YELLOW = "FFFF00"  # 黄色 = 新产品线
GREEN = "92D050"   # 绿色 = 成熟产品线

_A_NS = "{http://schemas.openxmlformats.org/drawingml/2006/main}"

# slide1 标题。
SLIDE1_TITLE = "ACME Hong Kong — Product Mix"
# slide1 含数字文本框（指标性叙述数字，glossary 命中 REVENUE）。
SLIDE1_TEXTBOX = "FY2024 REVENUE reached US$2,680m, up 14% YoY"
# slide2 演讲者备注（glossary 命中 PROFIT）。
SLIDE2_NOTES = "PROFIT for FY2024 was US$6,950m pending final audit"
# slide2 无数字文本框（不应被收进 note fragments，反例）。
SLIDE2_PLAIN_TEXTBOX = "Performance overview by reporting period"

# --- slide1 表格：行=产品线，列=FY2022..FY2024 -------------------------------
# 行 0 = 表头行（左上角空 + 三期间列）；行 1..3 = 产品线数据行。
# 填充策略：
#   Critical Illness 行（新产品线）-> 黄色 FFFF00。
#   Whole Life 行（成熟产品线）    -> 绿色 92D050。
#   Term Life 行                  -> 一格 theme 色（accent1）+ 其余无填充。
SLIDE1_COL_HEADERS = ["FY2022", "FY2023", "FY2024"]
SLIDE1_ROWS = [
    # (行标签, [三列数值字符串], 填充方案)
    ("Critical Illness", ["120", "180", "240"], "yellow"),
    ("Whole Life", ["980", "1010", "1055"], "green"),
    ("Term Life", ["310", "330", "355"], "term_mixed"),
]

# slide2 转置表：第一行 = ['Period', 'REVENUE', 'NEWSALES', 'PROFIT']，数据行=期间。
SLIDE2_HEADER = ["Period", "REVENUE", "NEWSALES", "PROFIT"]
SLIDE2_ROWS = [
    ("FY2023", ["2350", "4200", "2050"]),
    ("FY2024", ["2680", "4750", "2210"]),
]

# accent1 在默认 pptx 主题里的 srgb（self_verify 会从实际 theme 解析核对）。
ACCENT1_FALLBACK = "4F81BD"


# ---------------------------------------------------------------------------
# 主题色解析（从 ppt/theme1.xml 读 accent1，确保 ground truth 与实际写入一致）
# ---------------------------------------------------------------------------

def _resolve_accent1_rgb(prs: Presentation) -> str:
    """从 slide master 关联的 theme1.xml 解析 accent1 的真实 'RRGGBB'。"""
    master = prs.slide_masters[0]
    for rel in master.part.rels.values():
        if "theme" in rel.reltype:
            root = ET.fromstring(rel.target_part.blob)
            scheme = root.find(f".//{_A_NS}clrScheme")
            if scheme is None:
                break
            node = scheme.find(f"{_A_NS}accent1")
            if node is None:
                break
            srgb = node.find(f"{_A_NS}srgbClr")
            if srgb is not None and srgb.get("val"):
                return srgb.get("val").upper()
            sysc = node.find(f"{_A_NS}sysClr")
            if sysc is not None and sysc.get("lastClr"):
                return sysc.get("lastClr").upper()
    return ACCENT1_FALLBACK


# ---------------------------------------------------------------------------
# 构建幻灯片
# ---------------------------------------------------------------------------

def _set_cell_rgb(cell, rgb_hex: str) -> None:
    cell.fill.solid()
    cell.fill.fore_color.rgb = RGBColor.from_string(rgb_hex)


def _set_cell_theme(cell) -> None:
    cell.fill.solid()
    cell.fill.fore_color.theme_color = MSO_THEME_COLOR.ACCENT_1


def _build_slide1(prs: Presentation, accent1_rgb: str) -> dict:
    """slide1：标题 + 带填充色表格 + 含数字文本框。返回该页 ground truth。"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    # 标题文本框
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_box.text_frame.text = SLIDE1_TITLE

    # 表格：(1 表头 + 3 数据) 行 × (1 标签 + 3 期间) 列 = 4 × 4
    n_rows = 1 + len(SLIDE1_ROWS)
    n_cols = 1 + len(SLIDE1_COL_HEADERS)
    gtbl = slide.shapes.add_table(
        n_rows, n_cols, Inches(0.5), Inches(1.3), Inches(9), Inches(3)
    ).table

    # 逐格真值累积（cell_ref 'R{r}C{c}' 1-based）。
    cells_truth: dict[str, dict] = {}

    def _truth(r: int, c: int, value, rgb):
        ref = f"R{r}C{c}"
        cells_truth[ref] = {"cell_ref": ref, "value": value, "resolved_rgb": rgb}

    # 表头行（R1）：C1 空（左上角，无填充）、C2..C4 期间（无填充）
    gtbl.cell(0, 0).text = ""
    _truth(1, 1, "", None)
    for j, col in enumerate(SLIDE1_COL_HEADERS, start=1):
        gtbl.cell(0, j).text = col
        _truth(1, j + 1, col, None)

    # 数据行（R2..R4）
    for i, (label, vals, scheme) in enumerate(SLIDE1_ROWS, start=1):
        r = i + 1  # 1-based grid row
        gtbl.cell(i, 0).text = label
        if scheme == "yellow":
            label_rgb = YELLOW
            _set_cell_rgb(gtbl.cell(i, 0), YELLOW)
        elif scheme == "green":
            label_rgb = GREEN
            _set_cell_rgb(gtbl.cell(i, 0), GREEN)
        else:  # term_mixed：标签格无填充
            label_rgb = None
        _truth(r, 1, label, label_rgb)

        for j, v in enumerate(vals, start=1):
            c = j + 1  # 1-based grid col
            gtbl.cell(i, j).text = v
            if scheme == "yellow":
                _set_cell_rgb(gtbl.cell(i, j), YELLOW)
                _truth(r, c, v, YELLOW)
            elif scheme == "green":
                _set_cell_rgb(gtbl.cell(i, j), GREEN)
                _truth(r, c, v, GREEN)
            else:  # term_mixed：仅第一格（FY2022）用 theme 色，其余无填充
                if j == 1:
                    _set_cell_theme(gtbl.cell(i, j))
                    _truth(r, c, v, accent1_rgb)
                else:
                    _truth(r, c, v, None)

    # 含数字文本框
    box = slide.shapes.add_textbox(Inches(0.5), Inches(4.6), Inches(9), Inches(0.8))
    box.text_frame.text = SLIDE1_TEXTBOX

    return {
        "sheet": "slide1_table1",
        "n_rows": n_rows,
        "n_cols": n_cols,
        "cells": cells_truth,
    }


def _build_slide2(prs: Presentation) -> dict:
    """slide2：转置表（期间在行、指标在列）+ 含数字备注 + 无数字文本框。"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 无数字文本框（反例：不应进 note fragments）
    plain = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    plain.text_frame.text = SLIDE2_PLAIN_TEXTBOX

    # 转置表：4 行（1 表头 + 2 期间）× 4 列（标签 + 3 指标）
    n_rows = 1 + len(SLIDE2_ROWS)
    n_cols = len(SLIDE2_HEADER)
    tbl = slide.shapes.add_table(
        n_rows, n_cols, Inches(0.5), Inches(1.1), Inches(9), Inches(2.5)
    ).table

    cells_truth: dict[str, dict] = {}

    def _truth(r: int, c: int, value):
        ref = f"R{r}C{c}"
        cells_truth[ref] = {"cell_ref": ref, "value": value, "resolved_rgb": None}

    # 表头行（无填充）
    for j, h in enumerate(SLIDE2_HEADER):
        tbl.cell(0, j).text = h
        _truth(1, j + 1, h)
    # 数据行（无填充）
    for i, (period, vals) in enumerate(SLIDE2_ROWS, start=1):
        tbl.cell(i, 0).text = period
        _truth(i + 1, 1, period)
        for j, v in enumerate(vals, start=1):
            tbl.cell(i, j).text = v
            _truth(i + 1, j + 1, v)

    # 演讲者备注（含数字）
    slide.notes_slide.notes_text_frame.text = SLIDE2_NOTES

    return {
        "sheet": "slide2_table1",
        "n_rows": n_rows,
        "n_cols": n_cols,
        "cells": cells_truth,
        "orientation": "transposed",
    }


# ---------------------------------------------------------------------------
# note fragments 期望（含数字句段 + glossary 命中）
# ---------------------------------------------------------------------------

def _note_fragments_truth() -> list[dict]:
    """叙述层含数字句段期望（确定性规则：含 digit；glossary_hits 经词典命中）。

    slide1 文本框 'FY2024 REVENUE reached US$2,680m, up 14% YoY' -> 命中 REVENUE。
    slide2 演讲者备注 'PROFIT for FY2024 was US$6,950m pending final audit' -> 命中 PROFIT。
    slide2 无数字文本框（'Performance overview by reporting period'）不应出现。
    """
    return [
        {
            "slide_no": 1,
            "source_kind": "textbox",
            "text": SLIDE1_TEXTBOX,
            "glossary_hits": ["REVENUE"],
        },
        {
            "slide_no": 2,
            "source_kind": "notes",
            "text": SLIDE2_NOTES,
            "glossary_hits": ["PROFIT"],
        },
    ]


# ---------------------------------------------------------------------------
# OCR fake 测试向量（针对已存在的 scanned.pdf 三页）
# ---------------------------------------------------------------------------

# scanned.pdf 的三页都是 digital.pdf 表格页的渲染图，真值沿用其表格数值。
# 表格布局：R1 = 期间表头（C1 空 + FY2022/FY2023/FY2024）；R2..R4 = REVENUE/NEWSALES/PROFIT。
_OCR_TABLE_LAYOUT = {
    "R1C2": "FY2022", "R1C3": "FY2023", "R1C4": "FY2024",
    "R2C1": "REVENUE", "R2C2": "2100", "R2C3": "2350", "R2C4": "2680",
    "R3C1": "NEWSALES", "R3C2": "3800", "R3C3": "4200", "R3C4": "4750",
    "R4C1": "PROFIT", "R4C2": "1900", "R4C3": "2050", "R4C4": "2210",
}
# 低置信注入：每页同样两格 confidence < 0.85（其余 0.99）。
_OCR_LOW_CONF_REFS = {"R2C2", "R4C4"}  # REVENUE·FY2022 与 PROFIT·FY2024
_OCR_HIGH_CONF = 0.99
_OCR_LOW_CONF = 0.40
_OCR_MIN_CONFIDENCE = 0.85


def _ocr_page_result(page_no: int) -> dict:
    """单页 FakeBackend 应返回的 OcrPageResult 数据（一张 4×4 表 + 2 低置信格）。"""
    cells = []
    for ref, text in _OCR_TABLE_LAYOUT.items():
        r = int(ref[1])
        c = int(ref[3])
        conf = _OCR_LOW_CONF if ref in _OCR_LOW_CONF_REFS else _OCR_HIGH_CONF
        cells.append({"row": r, "col": c, "text": text, "confidence": conf})
    return {
        "page_no": page_no,
        "tables": [{"n_rows": 4, "n_cols": 4, "cells": cells}],
        "warnings": [],
    }


def _ocr_fake_vectors() -> dict:
    """OCR fake 测试向量：3 页输入 + 期望的网格 / 入队结果。

    scanned.pdf 有 3 页（二期 fixture 设计），每页一张 4×4 表、各混 2 个低置信格。
    期望：
        - 产出 3 张 StyledGrid（sheet 'page{N}_table1'）。
        - 每格 StyledCell.confidence = OCR 置信度、resolved_rgb=None。
        - 低于 min_confidence(0.85) 的格仍入网格，但每张 grid 加 warning，
          且给了 queue 时按 reason='low_confidence_ocr'、priority=30 入队。
        - 总入队数 = 3 页 × 2 低置信格 = 6。
    """
    pages = [_ocr_page_result(n) for n in (1, 2, 3)]
    expected_grids = []
    for n in (1, 2, 3):
        expected_grids.append({
            "sheet": f"page{n}_table1",
            "n_rows": 4,
            "n_cols": 4,
            "n_cells": len(_OCR_TABLE_LAYOUT),
            "low_confidence_refs": sorted(_OCR_LOW_CONF_REFS),
            "expect_warning": True,
        })
    return {
        "source_pdf": SCANNED_PDF_REL,
        "min_confidence": _OCR_MIN_CONFIDENCE,
        "low_confidence_value": _OCR_LOW_CONF,
        "high_confidence_value": _OCR_HIGH_CONF,
        "enqueue_reason": "low_confidence_ocr",
        "enqueue_priority": 30,
        "pages": pages,
        "expected_grids": expected_grids,
        "expected_total_enqueued": 3 * len(_OCR_LOW_CONF_REFS),
        "cell_values": dict(_OCR_TABLE_LAYOUT),
    }


# ---------------------------------------------------------------------------
# 组装 + 落盘 + 自校验
# ---------------------------------------------------------------------------

def build_deck() -> dict:
    prs = Presentation()
    accent1_rgb = _resolve_accent1_rgb(prs)

    slide1 = _build_slide1(prs, accent1_rgb)
    slide2 = _build_slide2(prs)

    OUT_DIR.mkdir(parents=True, exist_ok=True)
    prs.save(str(PPTX_PATH))

    ground_truth = {
        "file": PPTX_PATH.name,
        "theme": {"accent1_resolved_rgb": accent1_rgb},
        "slides": {
            "slide1": slide1,
            "slide2": slide2,
        },
        "note_fragments": _note_fragments_truth(),
        "ocr_fake": _ocr_fake_vectors(),
    }
    GT_PATH.write_text(
        json.dumps(ground_truth, ensure_ascii=False, indent=2), encoding="utf-8"
    )
    return ground_truth


def self_verify(gt: dict) -> None:
    """用 python-pptx 读回，断言值 / 填充 / 备注写入成功（fixture 可靠性自校验）。"""
    prs = Presentation(str(PPTX_PATH))
    slides = list(prs.slides)
    assert len(slides) == 2, len(slides)
    accent1_rgb = gt["theme"]["accent1_resolved_rgb"]

    # --- slide1：标题 + 表格填充 + 文本框 ---
    s1 = slides[0]
    texts1 = [sh.text_frame.text for sh in s1.shapes if sh.has_text_frame]
    assert SLIDE1_TITLE in texts1, texts1
    assert SLIDE1_TEXTBOX in texts1, texts1

    tables1 = [sh.table for sh in s1.shapes if sh.has_table]
    assert len(tables1) == 1, len(tables1)
    t1 = tables1[0]
    assert len(t1.rows) == 4 and len(t1.columns) == 4

    # 黄色行（Critical Illness, 数据格 R2C2..R2C4 = FFFF00）
    assert t1.cell(0, 1).text == "FY2022"
    assert t1.cell(1, 0).text == "Critical Illness"
    assert t1.cell(1, 1).fill.fore_color.type is not None
    assert str(t1.cell(1, 1).fill.fore_color.rgb) == YELLOW, t1.cell(1, 1).fill.fore_color.rgb
    # 绿色行（Whole Life）
    assert t1.cell(2, 0).text == "Whole Life"
    assert str(t1.cell(2, 1).fill.fore_color.rgb) == GREEN
    # theme 行（Term Life）：R4C2(=cell(3,1)) 是 theme 色 accent1
    assert t1.cell(3, 0).text == "Term Life"
    from pptx.enum.dml import MSO_FILL
    tc = t1.cell(3, 1).fill.fore_color
    assert tc.type == MSO_COLOR_TYPE.SCHEME, tc.type
    assert tc.theme_color == MSO_THEME_COLOR.ACCENT_1, tc.theme_color
    # 无填充格：Term Life 的 R4C3(=cell(3,2))；未显式设填充 -> fill.type 为 None
    # （继承表样式），抽取器须解析为 resolved_rgb=None。
    nofill = t1.cell(3, 2).fill
    assert nofill.type in (None, MSO_FILL.BACKGROUND), nofill.type

    # --- slide2：转置表 + 备注 + 无数字文本框 ---
    s2 = slides[1]
    texts2 = [sh.text_frame.text for sh in s2.shapes if sh.has_text_frame]
    assert SLIDE2_PLAIN_TEXTBOX in texts2, texts2
    tables2 = [sh.table for sh in s2.shapes if sh.has_table]
    assert len(tables2) == 1
    t2 = tables2[0]
    assert t2.cell(0, 0).text == "Period"
    assert t2.cell(0, 1).text == "REVENUE"
    assert t2.cell(1, 0).text == "FY2023"
    assert t2.cell(2, 3).text == "2210"
    # 演讲者备注
    assert s2.has_notes_slide
    assert SLIDE2_NOTES in s2.notes_slide.notes_text_frame.text

    # --- ground truth 自洽性 ---
    # theme 解析与 slide1 真值里 theme 格一致
    term_theme_ref = "R4C2"
    assert gt["slides"]["slide1"]["cells"][term_theme_ref]["resolved_rgb"] == accent1_rgb
    # note fragments：2 条、命中 REVENUE / PROFIT
    nf = gt["note_fragments"]
    assert {f["glossary_hits"][0] for f in nf} == {"REVENUE", "PROFIT"}, nf
    # OCR fake：3 页、每页 2 低置信、总入队 6
    ocr = gt["ocr_fake"]
    assert len(ocr["pages"]) == 3, len(ocr["pages"])
    assert ocr["expected_total_enqueued"] == 6
    for page in ocr["pages"]:
        lows = [c for c in page["tables"][0]["cells"] if c["confidence"] < ocr["min_confidence"]]
        assert len(lows) == 2, lows
    # scanned.pdf 必须已存在（OCR 向量引用它）
    assert SCANNED_PDF_PATH.exists(), f"缺少 {SCANNED_PDF_PATH}，请先跑 make_fixtures_pdf.py"

    print("self-verify: slide 值/填充/备注 + note fragments + OCR fake 向量 全部通过")


def main() -> None:
    gt = build_deck()
    self_verify(gt)
    n_cells = sum(len(s["cells"]) for s in gt["slides"].values())
    print(f"fixture 已生成于 {OUT_DIR}")
    print(f"  pptx: {PPTX_PATH.name}（2 slides, {n_cells} 逐格真值）")
    print(f"  ground_truth: {GT_PATH.name}")
    print(f"  theme accent1 -> {gt['theme']['accent1_resolved_rgb']}")
    print(f"  note fragments: {len(gt['note_fragments'])} 条")
    print(f"  OCR fake 向量: {len(gt['ocr_fake']['pages'])} 页, 期望入队 {gt['ocr_fake']['expected_total_enqueued']}")


if __name__ == "__main__":
    main()
