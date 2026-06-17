"""PPT 增强抽取契约（三期 PPT 增强，story #12/#13）—— 新模块，与旧 pptx_extractor 并存。

旧 ragspine/extraction/extractors/pptx_extractor.py 是 MVP 确定性抽取器（表格 + 原生图表 -> Fact），
38 个 e2e 测试依赖它、**不许改**。本模块是三期的「样式感知 + 叙述数字」增强，
产出统一中间表示 StyledGrid（与 xlsx/pdf 抽取器对齐），并新增叙述层数字抽取
（文本框 + 演讲者备注），让藏在叙述里的关键数字不会漏掉（PRD user stories 12、13）。

两条抽取路径：
    1) extract_grids        —— 原生表格 -> 样式感知网格（含填充色解析，走与 Excel
                               同一套颜色映射资产，user story 13）。
    2) extract_note_fragments —— 文本框 + 演讲者备注里**含数字**的句段 -> NoteFragment
                               （确定性规则：含 digit；glossary 命中指标代码，user story 12）。

命名与坐标约定（与 pdf_digital_extractor 对齐）：
    - sheet 命名 'slide{N}_table{M}'（N=幻灯片号 1-based，M=该页内表序 1-based）。
    - cell_ref 用 'R{行}C{列}'（1-based），表明来源是表格而非工作表坐标。
    - value = 空白归一化后的单元格文本（字符串，不做类型推断）。
    - resolved_rgb = 单元格填充色解析为 'RRGGBB' 大写十六进制；无填充 / 无法解析 -> None。

填充色解析约定（python-pptx）：
    - solid fill 的显式 RGB（MSO_FILL.SOLID + fore_color.type == RGB）直取。
    - theme 色（fore_color.type == SCHEME）经 ppt/theme1.xml 的 clrScheme 解析为真实
      RGB（思路参考 xlsx_styled_extractor 的 theme 解析）。
    - 无填充 / 继承 / 无法解析 -> None，并视情况在 grid 级追加 warning。

不规则表（转置 / 多级表头）照常产出网格，不做语义判断（语义归 ingestion 层）；
但 n_rows / n_cols 必须如实反映表格真实维度。
"""

import colorsys
import hashlib
import re
import xml.etree.ElementTree as ET
from dataclasses import dataclass, field
from pathlib import Path

from pptx import Presentation
from pptx.enum.dml import MSO_COLOR_TYPE, MSO_FILL

from ragspine.common.glossary import METRIC_SYNONYMS
from ragspine.extraction.ir import StyledCell, StyledGrid

# 抽取器版本标识（写入血缘，便于回归门禁区分解析器版本）。
EXTRACTOR_VERSION = "pptx_styled_v0"

# NoteFragment.source_kind 取值。
SOURCE_TEXTBOX = "textbox"   # 幻灯片正文文本框。
SOURCE_NOTES = "notes"       # 演讲者备注（notes slide）。

_A_NS = "{http://schemas.openxmlformats.org/drawingml/2006/main}"

# MSO_THEME_COLOR 枚举值 -> OOXML clrScheme 槽位名（pptx 主题槽位映射）。
# 注意 pptx 的 dk1/lt1 与 BACKGROUND/TEXT 经 clrMap 映射，但显式 accentN 直接对应
# clrScheme 的同名子元素，是本线（accent1）的关键路径。
_THEME_SLOT_BY_ENUM: dict[int, str] = {
    1: "dk1",       # DARK_1
    2: "lt1",       # LIGHT_1
    3: "dk2",       # DARK_2
    4: "lt2",       # LIGHT_2
    5: "accent1",   # ACCENT_1
    6: "accent2",   # ACCENT_2
    7: "accent3",   # ACCENT_3
    8: "accent4",   # ACCENT_4
    9: "accent5",   # ACCENT_5
    10: "accent6",  # ACCENT_6
    11: "hlink",    # HYPERLINK
    12: "folHlink", # FOLLOWED_HYPERLINK
    13: "dk1",      # TEXT_1（默认 clrMap：tx1 -> dk1）
    14: "lt1",      # BACKGROUND_1（默认 clrMap：bg1 -> lt1）
    15: "dk2",      # TEXT_2（默认 clrMap：tx2 -> dk2）
    16: "lt2",      # BACKGROUND_2（默认 clrMap：bg2 -> lt2）
}


@dataclass
class NoteFragment:
    """PPT 叙述层里一个含数字的句段（PRD user story 12）。

    字段语义约定：
        slide_no:       所在幻灯片号（1-based）。
        source_kind:    来源类型 'textbox'（正文文本框）/ 'notes'（演讲者备注）。
        locator:        原文精确定位回链（如 'slide2/notes' / 'slide1/textbox3'），
                        供 citation 与复核回指。
        text:           句段文本，空白归一化（首尾 strip + 内部连续空白折叠为单空格）。
        glossary_hits:  经 ragspine/common/glossary.py 命中的指标代码列表（确定性词典匹配，
                        如 ['REVENUE']；无命中为空列表，绝不做 LLM 推断）。
    """

    slide_no: int
    source_kind: str
    locator: str
    text: str
    glossary_hits: list[str] = field(default_factory=list)


def compute_file_hash(path: str | Path) -> str:
    """计算源文件内容 hash（十六进制串），作为版本血缘标识。"""
    digest = hashlib.sha256()
    with open(path, "rb") as f:
        for chunk in iter(lambda: f.read(65536), b""):
            digest.update(chunk)
    return digest.hexdigest()


def _normalize_text(text: object) -> str:
    """空白归一化：首尾 strip + 内部连续空白折叠为单空格。"""
    return " ".join(str(text).split())


# ---------------------------------------------------------------------------
# 主题色解析（从 slide 关联的 theme1.xml 的 clrScheme 读取真实 RGB）
# ---------------------------------------------------------------------------

def _scheme_slot_rgb(scheme: ET.Element, slot: str) -> str | None:
    """从 clrScheme 取某槽位的 'RRGGBB' 大写十六进制；无则 None。"""
    node = scheme.find(f"{_A_NS}{slot}")
    if node is None:
        return None
    srgb = node.find(f"{_A_NS}srgbClr")
    if srgb is not None and srgb.get("val"):
        return srgb.get("val").upper()
    sysc = node.find(f"{_A_NS}sysClr")
    if sysc is not None and sysc.get("lastClr"):
        return sysc.get("lastClr").upper()
    return None


def _theme_scheme_for_slide(slide) -> ET.Element | None:
    """沿 slide -> slide_layout -> slide_master -> theme 解析出 clrScheme 元素。"""
    try:
        master = slide.slide_layout.slide_master
    except (AttributeError, KeyError):
        return None
    for rel in master.part.rels.values():
        if rel.is_external:
            continue
        if "theme" in rel.reltype:
            root = ET.fromstring(rel.target_part.blob)
            return root.find(f".//{_A_NS}clrScheme")
    return None


def _apply_brightness(base_rgb: str, brightness: float | None) -> str:
    """按 python-pptx brightness（[-1, 1]）调整明度：>0 向白、<0 向黑。"""
    if not brightness:
        return base_rgb
    r = int(base_rgb[0:2], 16) / 255.0
    g = int(base_rgb[2:4], 16) / 255.0
    b = int(base_rgb[4:6], 16) / 255.0
    h, lum, s = colorsys.rgb_to_hls(r, g, b)
    if brightness < 0:
        lum = lum * (1.0 + brightness)
    else:
        lum = lum + (1.0 - lum) * brightness
    r2, g2, b2 = colorsys.hls_to_rgb(h, lum, s)
    return "{:02X}{:02X}{:02X}".format(
        round(r2 * 255), round(g2 * 255), round(b2 * 255)
    )


def _resolve_cell_fill(cell, scheme: ET.Element | None) -> str | None:
    """解析表格单元格填充色为真实 'RRGGBB'；无填充 / 无法解析返回 None。"""
    fill = cell.fill
    try:
        fill_type = fill.type
    except (TypeError, KeyError):
        return None
    if fill_type != MSO_FILL.SOLID:
        return None
    fore = fill.fore_color
    try:
        color_type = fore.type
    except (TypeError, KeyError):
        return None
    if color_type == MSO_COLOR_TYPE.RGB:
        return str(fore.rgb).upper()
    if color_type == MSO_COLOR_TYPE.SCHEME:
        if scheme is None:
            return None
        slot = _THEME_SLOT_BY_ENUM.get(int(fore.theme_color))
        if slot is None:
            return None
        base = _scheme_slot_rgb(scheme, slot)
        if base is None:
            return None
        try:
            brightness = float(fore.brightness)
        except (TypeError, ValueError):
            brightness = 0.0
        return _apply_brightness(base, brightness)
    return None


def _build_grid(slide, slide_no: int, table_no: int, table, doc_id: str,
                file_hash: str, scheme: ET.Element | None) -> StyledGrid:
    """把一张原生表格 shape 构建为 StyledGrid。"""
    n_rows = len(table.rows)
    n_cols = len(table.columns)
    grid = StyledGrid(
        sheet=f"slide{slide_no}_table{table_no}",
        source_doc_id=doc_id,
        source_file_hash=file_hash,
        n_rows=n_rows,
        n_cols=n_cols,
    )
    for r in range(n_rows):
        for c in range(n_cols):
            cell = table.cell(r, c)
            ref = f"R{r + 1}C{c + 1}"
            grid.cells[ref] = StyledCell(
                value=_normalize_text(cell.text),
                cell_ref=ref,
                resolved_rgb=_resolve_cell_fill(cell, scheme),
                confidence=None,
            )
    return grid


def extract_grids(path: str | Path) -> list[StyledGrid]:
    """抽取一个 pptx 的全部原生表格 -> list[StyledGrid]（每张表一个 StyledGrid）。

    每个 StyledGrid：
        sheet            = 'slide{N}_table{M}'（N 幻灯片号、M 该页表序，均 1-based）。
        source_doc_id    = 文件名；source_file_hash = 内容 hash。
        cells            = 'R{行}C{列}' -> StyledCell；value 为空白归一化单元格文本，
                           resolved_rgb 为填充色解析结果（显式 RGB 直取 / theme 色经
                           ppt/theme1.xml 解析 / 无填充 None）。
        n_rows / n_cols  = 表格真实逻辑行列数（不规则表也如实，语义判断归下游）。

    无表格 -> 返回 []。无法解析的填充色 -> 该格 resolved_rgb=None 并视情况加 grid warning。
    """
    path = Path(path)
    doc_id = path.name
    file_hash = compute_file_hash(path)
    prs = Presentation(str(path))

    grids: list[StyledGrid] = []
    for slide_no, slide in enumerate(prs.slides, start=1):
        scheme = _theme_scheme_for_slide(slide)
        table_no = 0
        for shape in slide.shapes:
            if not shape.has_table:
                continue
            table_no += 1
            grids.append(
                _build_grid(slide, slide_no, table_no, shape.table,
                            doc_id, file_hash, scheme)
            )
    return grids


# ---------------------------------------------------------------------------
# 叙述层含数字句段抽取（文本框 + 演讲者备注）
# ---------------------------------------------------------------------------

def _build_glossary_matchers() -> list[tuple[str, re.Pattern[str]]]:
    """为每个指标同义词构建 (metric_code, 大小写不敏感词边界正则)。

    仅用 METRIC_SYNONYMS（指标代码），不含 entity/geography，避免臆造其它指标。
    ASCII 同义词用 \\b 词边界匹配；含非 ASCII（中文）同义词退化为直接子串包含。
    """
    matchers: list[tuple[str, re.Pattern[str]]] = []
    for synonym, code in METRIC_SYNONYMS.items():
        if synonym.isascii():
            pattern = re.compile(rf"\b{re.escape(synonym)}\b", re.IGNORECASE)
        else:
            pattern = re.compile(re.escape(synonym))
        matchers.append((code, pattern))
    return matchers


_GLOSSARY_MATCHERS = _build_glossary_matchers()


def _glossary_hits(text: str) -> list[str]:
    """确定性词典命中：返回文本命中的指标代码（去重、稳定顺序）。"""
    hits: list[str] = []
    for code, pattern in _GLOSSARY_MATCHERS:
        if code in hits:
            continue
        if pattern.search(text):
            hits.append(code)
    return hits


def _has_digit(text: str) -> bool:
    return any(ch.isdigit() for ch in text)


def _iter_textbox_texts(slide):
    """产出该 slide 内非表格文本框的归一化文本（按 shape 出现顺序）。"""
    for shape in slide.shapes:
        if shape.has_table:
            continue
        if not shape.has_text_frame:
            continue
        text = _normalize_text(shape.text_frame.text)
        if text:
            yield text


def _notes_text(slide) -> str | None:
    """该 slide 演讲者备注的归一化文本；无备注返回 None。"""
    if not slide.has_notes_slide:
        return None
    text = _normalize_text(slide.notes_slide.notes_text_frame.text)
    return text or None


def extract_note_fragments(path: str | Path) -> list[NoteFragment]:
    """抽取一个 pptx 各幻灯片文本框 + 演讲者备注中**含数字**的句段。

    确定性规则（零 LLM）：
        - 遍历每页所有文本框（source_kind='textbox'）与演讲者备注（source_kind='notes'）。
        - 只收集**含 digit** 的 run / 段落（叙述里的指标性数字线索）。
        - text 空白归一化。
        - glossary_hits = 经 ragspine/common/glossary.py 确定性词典命中的指标代码列表
          （如句段含 'REVENUE' -> ['REVENUE']；'PROFIT' -> ['PROFIT']）。

    返回顺序按幻灯片号、再按页内出现顺序。无含数字句段 -> 返回 []。
    """
    path = Path(path)
    prs = Presentation(str(path))

    fragments: list[NoteFragment] = []
    for slide_no, slide in enumerate(prs.slides, start=1):
        box_no = 0
        for text in _iter_textbox_texts(slide):
            box_no += 1
            if not _has_digit(text):
                continue
            fragments.append(NoteFragment(
                slide_no=slide_no,
                source_kind=SOURCE_TEXTBOX,
                locator=f"slide{slide_no}/textbox{box_no}",
                text=text,
                glossary_hits=_glossary_hits(text),
            ))

        notes = _notes_text(slide)
        if notes is not None and _has_digit(notes):
            fragments.append(NoteFragment(
                slide_no=slide_no,
                source_kind=SOURCE_NOTES,
                locator=f"slide{slide_no}/notes",
                text=notes,
                glossary_hits=_glossary_hits(notes),
            ))

    fragments.sort(key=lambda f: f.slide_no)
    return fragments
