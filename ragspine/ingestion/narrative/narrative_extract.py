"""叙事文本抽取：真实文档 -> 可直接喂 ragspine/retrieval/chunking/chunking.py 的「文档文本 + 定位段」。

叙事通路（Narrative Path）吃的是归因 / 监管 / 会议纪要类文本（docs/architecture.md「Channels」），
本模块只做确定性文本抽取，零 OCR、零 LLM：

    - pptx：python-pptx 遍历每页文本框（含占位符）+ 演讲者备注，按 slide 聚合；
      表格内容跳过（表格数字归结构化通路管）。定位串 'slide={N},frame={M}'
      （M 只给产出非空文本的非表格文本框编号，1-based）/ 'slide={N},notes'。
    - PDF（数字型）：pypdfium2 textpage 逐页取文本，定位串 'page={N}'（真实页号）。
      扫描型页（无文本层）本期不做（OCR 线另有归属）：跳过、计数并告警。

统一返回 NarrativeDoc：segments 各段文本内部以 '\\n' 分行、to_text() 以空行
（'\\n\\n'）连接各段，正好配合 chunking 的「非空行 = 段落」切分契约。
"""

from dataclasses import dataclass, field
from pathlib import Path

import pypdfium2 as pdfium
from pptx import Presentation

from ragspine.extraction.extractors.pptx_styled_extractor import compute_file_hash

# 本期支持的叙事来源后缀（扫描件 / 其它格式归别的线）。
SUPPORTED_SUFFIXES = {".pptx", ".pdf"}


@dataclass
class NarrativeSegment:
    """一个叙事文本段：文本 + 原文定位串。

    source_locator 取值约定：
        'slide={N},frame={M}'  pptx 第 N 页第 M 个产出文本的文本框（均 1-based）。
        'slide={N},notes'      pptx 第 N 页演讲者备注。
        'page={N}'             PDF 第 N 页（1-based，真实页号，跳过页不占用）。
    """

    text: str
    source_locator: str


@dataclass
class NarrativeDoc:
    """一个文档的叙事抽取结果（喂 chunking 的统一中间表示）。

    字段语义约定：
        doc_id:        源文件名（与全库血缘约定一致）。
        file_hash:     源文件内容 sha256 十六进制串（版本血缘）。
        segments:      叙事文本段列表（按 slide/页序）。
        skipped_pages: 无文本层被跳过的 PDF 页数（pptx 恒为 0）。
        warnings:      告警汇聚（如逐页跳过原因）。
    """

    doc_id: str
    file_hash: str
    segments: list[NarrativeSegment] = field(default_factory=list)
    skipped_pages: int = 0
    warnings: list[str] = field(default_factory=list)

    def to_text(self) -> str:
        """各段以空行连接成文档级纯文本（chunking 的输入契约）。"""
        return "\n\n".join(seg.text for seg in self.segments)


def _clean_block(text: str) -> str:
    """块内文本归一化：逐行折叠内部空白、丢空行，行间以 '\\n' 连接。

    pptx 软换行（'\\v'）视同换行；返回空串表示该块无实质文本。
    """
    lines = (" ".join(line.split()) for line in text.replace("\v", "\n").splitlines())
    return "\n".join(line for line in lines if line)


def extract_pptx_narrative(path: str | Path) -> NarrativeDoc:
    """抽取一个 pptx 的叙事文本：每页文本框（含占位符）+ 演讲者备注。

    - 表格 shape 跳过（表格数字归结构化通路）。
    - 仅产出非空文本的文本框获得 frame 序号（与 pptx_styled_extractor 的
      box_no 约定一致）；备注段排在该页文本框段之后。
    """
    path = Path(path)
    doc = NarrativeDoc(doc_id=path.name, file_hash=compute_file_hash(path))
    prs = Presentation(str(path))

    for slide_no, slide in enumerate(prs.slides, start=1):
        frame_no = 0
        for shape in slide.shapes:
            if shape.has_table or not shape.has_text_frame:
                continue
            text = _clean_block(shape.text_frame.text)
            if not text:
                continue
            frame_no += 1
            doc.segments.append(NarrativeSegment(
                text=text,
                source_locator=f"slide={slide_no},frame={frame_no}",
            ))
        if slide.has_notes_slide:
            notes = _clean_block(slide.notes_slide.notes_text_frame.text)
            if notes:
                doc.segments.append(NarrativeSegment(
                    text=notes,
                    source_locator=f"slide={slide_no},notes",
                ))
    return doc


def extract_pdf_narrative(path: str | Path) -> NarrativeDoc:
    """抽取一个数字型 PDF 的叙事文本：pypdfium2 textpage 逐页取文本按页聚合。

    无文本层的页（扫描型形态）跳过：skipped_pages 计数 + warnings 点名页号，
    OCR 归扫描线处理，本模块绝不臆造内容。
    """
    path = Path(path)
    doc = NarrativeDoc(doc_id=path.name, file_hash=compute_file_hash(path))
    pdf = pdfium.PdfDocument(str(path))
    try:
        for idx in range(len(pdf)):
            page = pdf[idx]
            try:
                textpage = page.get_textpage()
                try:
                    raw = textpage.get_text_range()
                finally:
                    textpage.close()
            finally:
                page.close()
            text = _clean_block(raw)
            if text:
                doc.segments.append(NarrativeSegment(
                    text=text,
                    source_locator=f"page={idx + 1}",
                ))
            else:
                doc.skipped_pages += 1
                doc.warnings.append(
                    f"page={idx + 1}: 无文本层，跳过（扫描页归 OCR 线处理）"
                )
    finally:
        pdf.close()
    return doc


def extract_narrative(path: str | Path) -> NarrativeDoc:
    """按后缀分发到对应抽取器；不支持的后缀 ValueError。"""
    path = Path(path)
    suffix = path.suffix.lower()
    if suffix == ".pptx":
        return extract_pptx_narrative(path)
    if suffix == ".pdf":
        return extract_pdf_narrative(path)
    raise ValueError(f"不支持的叙事来源类型：{path.name}（仅 .pptx / .pdf）")
