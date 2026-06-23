"""生成合成数据：pptx（表格页×2 + 原生图表页×1）+ xlsx（5yr）+ ground_truth.json。

数值全部硬编码、确定性，作为端到端验证的唯一真值来源。
"""

import json
from pathlib import Path
from typing import Any

import rootutils
from openpyxl import Workbook
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.presentation import Presentation as PresentationDoc
from pptx.util import Inches

ROOT_DIR = rootutils.find_root(Path(__file__), indicator=".project-root")

OUT_DIR = ROOT_DIR / "data" / "synthetic"
PPTX_PATH = OUT_DIR / "ACME_FY2024_Review.pptx"
XLSX_PATH = OUT_DIR / "fiveyr_summary.xlsx"
GT_PATH = OUT_DIR / "ground_truth.json"

# --- 硬编码真值（单位：USD_M / PCT）---------------------------------------

# 香港表格页：REVENUE/NEWSALES/PROFIT/ROE × FY2022/FY2023/FY2024
HK_TABLE = {
    "REVENUE": {"FY2022": 2100.0, "FY2023": 2350.0, "FY2024": 2680.0},
    "NEWSALES": {"FY2022": 3800.0, "FY2023": 4200.0, "FY2024": 4750.0},
    "PROFIT": {"FY2022": 1900.0, "FY2023": 2050.0, "FY2024": 2210.0},
    "ROE": {"FY2022": 12.5, "FY2023": 13.1, "FY2024": 13.8},
}

# 中国表格页：不含 ROE（用于 ROE/ACME_CN/FY2024 的 not_found 测试）
CN_TABLE = {
    "REVENUE": {"FY2022": 1050.0, "FY2023": 1180.0, "FY2024": 1320.0},
    "NEWSALES": {"FY2022": 2200.0, "FY2023": 2500.0, "FY2024": 2900.0},
    "PROFIT": {"FY2022": 980.0, "FY2023": 1090.0, "FY2024": 1240.0},
}

# 集团 PROFIT 原生柱状图：FY2022-24
GROUP_PROFIT_CHART = {"FY2022": 5800.0, "FY2023": 6300.0, "FY2024": 6950.0}

# xlsx 5yr：ACME_HK 的 REVENUE/NEWSALES，FY2020-FY2024 + 2024H1
XLSX_PERIODS = ["FY2020", "FY2021", "FY2022", "FY2023", "FY2024", "2024H1"]
XLSX_HK = {
    "REVENUE": [1750.0, 1820.0, 2100.0, 2350.0, 2680.0, 1390.0],
    "NEWSALES": [3100.0, 3300.0, 3800.0, 4200.0, 4750.0, 2480.0],
}

# 期间字符串 -> (period_type, period)，与 glossary.normalize_period 一致
_PERIOD_MAP = {
    "FY2020": ("FY", "2020"),
    "FY2021": ("FY", "2021"),
    "FY2022": ("FY", "2022"),
    "FY2023": ("FY", "2023"),
    "FY2024": ("FY", "2024"),
    "2024H1": ("HY", "2024H1"),
}
_UNIT = {"REVENUE": "USD_M", "NEWSALES": "USD_M", "PROFIT": "USD_M", "ROE": "PCT"}


def _add_table_slide(
    prs: PresentationDoc,
    title_text: str,
    data: dict[str, dict[str, float]],
    periods: list[str],
) -> None:
    """添加一张"指标×期间"表格页。"""
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title_text

    metrics = list(data.keys())
    rows = len(metrics) + 1
    cols = len(periods) + 1
    table = slide.shapes.add_table(
        rows, cols, Inches(0.5), Inches(1.6), Inches(9.0), Inches(0.4 * rows)
    ).table

    table.cell(0, 0).text = "Metric"
    for j, period in enumerate(periods, start=1):
        table.cell(0, j).text = period
    for i, metric in enumerate(metrics, start=1):
        table.cell(i, 0).text = metric
        for j, period in enumerate(periods, start=1):
            table.cell(i, j).text = f"{data[metric][period]:,.1f}"


def _add_chart_slide(
    prs: PresentationDoc,
    title_text: str,
    series_name: str,
    chart_data: dict[str, float],
) -> None:
    """添加一张原生柱状图页（CategoryChartData，数值内嵌进图表 XML）。"""
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title_text

    cats = list(chart_data.keys())
    cd: Any = CategoryChartData()  # type: ignore[no-untyped-call]
    cd.categories = cats
    cd.add_series(series_name, [chart_data[c] for c in cats])
    slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(1.0),
        Inches(1.6),
        Inches(8.0),
        Inches(4.5),
        cd,
    )


def make_pptx() -> None:
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    _add_table_slide(
        prs, "ACME Hong Kong — Financial Performance (US$m)", HK_TABLE,
        ["FY2022", "FY2023", "FY2024"],
    )
    _add_table_slide(
        prs, "ACME China — Financial Performance (US$m)", CN_TABLE,
        ["FY2022", "FY2023", "FY2024"],
    )
    _add_chart_slide(prs, "ACME Group — PROFIT Trend", "PROFIT", GROUP_PROFIT_CHART)
    prs.save(str(PPTX_PATH))


def make_xlsx() -> None:
    wb = Workbook()
    ws = wb.active
    ws.title = "5yr"
    ws["A1"] = "ACME Hong Kong"  # A1 约定实体
    for j, period in enumerate(XLSX_PERIODS, start=2):
        ws.cell(row=1, column=j, value=period)
    for i, (metric, values) in enumerate(XLSX_HK.items(), start=2):
        ws.cell(row=i, column=1, value=metric)
        for j, val in enumerate(values, start=2):
            ws.cell(row=i, column=j, value=val)
    wb.save(str(XLSX_PATH))


def make_ground_truth() -> list[dict[str, object]]:
    """把所有写入数值平铺成真值清单。"""
    gt: list[dict[str, object]] = []

    def add(
        metric: str,
        entity: str,
        period_label: str,
        value: float,
        source: str,
    ) -> None:
        ptype, period = _PERIOD_MAP[period_label]
        gt.append({
            "metric": metric,
            "entity": entity,
            "period_type": ptype,
            "period": period,
            "value": value,
            "unit": _UNIT[metric],
            "source": source,
        })

    for metric, by_period in HK_TABLE.items():
        for period_label, value in by_period.items():
            add(metric, "ACME_HK", period_label, value, "pptx_table")
    for metric, by_period in CN_TABLE.items():
        for period_label, value in by_period.items():
            add(metric, "ACME_CN", period_label, value, "pptx_table")
    for period_label, value in GROUP_PROFIT_CHART.items():
        add("PROFIT", "ACME_GROUP", period_label, value, "pptx_chart")
    for metric, values in XLSX_HK.items():
        for period_label, value in zip(XLSX_PERIODS, values, strict=False):
            add(metric, "ACME_HK", period_label, value, "xlsx")

    GT_PATH.write_text(json.dumps(gt, ensure_ascii=False, indent=2), encoding="utf-8")
    return gt


def main() -> None:
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    make_pptx()
    make_xlsx()
    gt = make_ground_truth()
    print(f"合成数据已生成于 {OUT_DIR}")
    print(f"  pptx: {PPTX_PATH.name}")
    print(f"  xlsx: {XLSX_PATH.name}")
    print(f"  ground_truth: {GT_PATH.name}（{len(gt)} 条真值）")


if __name__ == "__main__":
    main()
