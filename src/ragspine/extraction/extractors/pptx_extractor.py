"""PPTX 确定性抽取：从原生表格与原生图表读出数值，零 OCR、零 LLM。

约定：
    - slide 标题（或首个文本框）经 glossary 解析出 entity / geography。
    - 表格：第一行=期间表头，第一列=指标名，行列都经 glossary 归一；
      无法识别的行/列跳过并记入 warnings。
    - 原生图表：series.name 经 glossary 归一为指标，categories 为期间，
      series.values 为数值（指标由 series 决定，实体由 slide 标题决定）。
      这是"PPT 图表内嵌真实数据"路径的证明点——数据来自 chart 的 XML，不靠图像。
所有渠道默认 'TOTAL'。
"""

from pathlib import Path

from pptx import Presentation
from pptx.chart.chart import Chart
from pptx.slide import Slide
from pptx.table import Table

from ragspine.common.glossary import (
    geography_for_entity,
    normalize_entity,
    normalize_metric,
    normalize_period,
    unit_for_metric,
)
from ragspine.storage.fact_store import Fact


def _slide_title(slide: Slide) -> str | None:
    """取 slide 标题占位符文本；没有则退回首个有文字的文本框。"""
    if slide.shapes.title is not None and slide.shapes.title.has_text_frame:
        text = str(slide.shapes.title.text_frame.text.strip())
        if text:
            return text
    for shape in slide.shapes:
        # pptx 静态类型在 BaseShape 上未暴露 text_frame，由 has_text_frame 守卫。
        if shape.has_text_frame and shape.text_frame.text.strip():  # type: ignore[attr-defined]
            return str(shape.text_frame.text.strip())  # type: ignore[attr-defined]
    return None


def _coerce_number(raw: str) -> float | None:
    """单元格文本转数值：去千分位逗号/货币符/百分号，失败返回 None。"""
    if raw is None:
        return None
    cleaned = raw.strip().replace(",", "").replace("$", "").replace("%", "")
    if not cleaned:
        return None
    try:
        return float(cleaned)
    except ValueError:
        return None


def _extract_table(
    table: Table,
    slide_idx: int,
    table_idx: int,
    entity: str,
    geography: str,
    doc_id: str,
    warnings: list[str],
) -> list[Fact]:
    facts: list[Fact] = []
    rows = list(table.rows)
    if not rows:
        return facts

    # 第一行：期间表头（首格为指标列标题，跳过）
    header = [cell.text for cell in rows[0].cells]
    periods: list[tuple[int, tuple[str, str], str] | None] = [None]
    for col_idx in range(1, len(header)):
        parsed = normalize_period(header[col_idx])
        if parsed is None:
            warnings.append(
                f"slide={slide_idx},table={table_idx}: 无法识别期间表头 '{header[col_idx]}'，跳过该列"
            )
            periods.append(None)
        else:
            periods.append((col_idx, parsed, header[col_idx].strip()))

    # 数据行：首列为指标名
    for row in rows[1:]:
        cells = [cell.text for cell in row.cells]
        metric_label = cells[0] if cells else ""
        metric_code = normalize_metric(metric_label)
        if metric_code is None:
            warnings.append(
                f"slide={slide_idx},table={table_idx}: 无法识别指标 '{metric_label}'，跳过该行"
            )
            continue
        unit = unit_for_metric(metric_code) or "USD_M"
        for col_idx in range(1, len(cells)):
            pinfo = periods[col_idx] if col_idx < len(periods) else None
            if pinfo is None:
                continue
            _, (period_type, period), period_raw = pinfo
            value = _coerce_number(cells[col_idx])
            if value is None:
                warnings.append(
                    f"slide={slide_idx},table={table_idx},row={metric_code},col={period_raw}: 空/非数值，跳过"
                )
                continue
            facts.append(
                Fact(
                    metric_code=metric_code,
                    entity=entity,
                    geography=geography,
                    channel="TOTAL",
                    period_type=period_type,
                    period=period,
                    value=value,
                    unit=unit,
                    source_doc_id=doc_id,
                    source_locator=f"slide={slide_idx},table={table_idx},row={metric_code},col={period_raw}",
                )
            )
    return facts


def _extract_chart(
    chart: Chart,
    slide_idx: int,
    chart_idx: int,
    entity: str,
    geography: str,
    doc_id: str,
    warnings: list[str],
) -> list[Fact]:
    facts: list[Fact] = []
    plot = chart.plots[0] if chart.plots else None
    if plot is None:
        return facts
    categories = [str(c) for c in plot.categories]
    parsed_cats: list[tuple[str, str] | None] = []
    for cat in categories:
        p = normalize_period(cat)
        if p is None:
            warnings.append(
                f"slide={slide_idx},chart={chart_idx}: 无法识别图表类目 '{cat}'，跳过"
            )
        parsed_cats.append(p)

    for series in plot.series:
        metric_code = normalize_metric(series.name)
        if metric_code is None:
            warnings.append(
                f"slide={slide_idx},chart={chart_idx}: 无法识别 series '{series.name}'，跳过"
            )
            continue
        unit = unit_for_metric(metric_code) or "USD_M"
        for cat_raw, parsed, value in zip(
            categories, parsed_cats, series.values, strict=False
        ):
            if parsed is None or value is None:
                continue
            period_type, period = parsed
            facts.append(
                Fact(
                    metric_code=metric_code,
                    entity=entity,
                    geography=geography,
                    channel="TOTAL",
                    period_type=period_type,
                    period=period,
                    value=float(value),
                    unit=unit,
                    source_doc_id=doc_id,
                    source_locator=f"slide={slide_idx},chart={chart_idx},series={metric_code},cat={cat_raw}",
                )
            )
    return facts


def extract_facts(path: str | Path) -> tuple[list[Fact], list[str]]:
    """抽取一个 pptx 的全部 Fact。返回 (facts, warnings)。"""
    path = Path(path)
    doc_id = path.name
    prs = Presentation(str(path))
    facts: list[Fact] = []
    warnings: list[str] = []

    for slide_idx, slide in enumerate(prs.slides, start=1):
        title = _slide_title(slide)
        entity = normalize_entity(title) if title else None
        if entity is None:
            warnings.append(f"slide={slide_idx}: 标题 '{title}' 无法解析实体，跳过该页")
            continue
        geography = geography_for_entity(entity) or "UNKNOWN"

        table_idx = 0
        chart_idx = 0
        for shape in slide.shapes:
            if shape.has_table:
                table_idx += 1
                facts.extend(
                    _extract_table(
                        shape.table, slide_idx, table_idx, entity, geography, doc_id, warnings
                    )
                )
            elif shape.has_chart:
                chart_idx += 1
                facts.extend(
                    _extract_chart(
                        shape.chart, slide_idx, chart_idx, entity, geography, doc_id, warnings
                    )
                )

    return facts, warnings
